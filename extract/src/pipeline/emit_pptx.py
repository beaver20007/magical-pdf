"""DocumentIR → editable PPTX (1 PDF page = 1 slide, all text editable)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from src.pipeline.cleanup import clean_ocr_text
from src.pipeline.gap_fill import _is_gap_path
from src.pipeline.ir import Bbox, DocumentIR, ImageBlock, Page, TableBlock, TextBlock
from src.pipeline.text_mask import mask_page_background, mask_text_regions

_PT_TO_EMU = 12700


def _page_size(page: Page | None, fallback_w: float, fallback_h: float) -> tuple[int, int]:
    w_pt = page.width_pt if page else fallback_w
    h_pt = page.height_pt if page else fallback_h
    return int(w_pt * _PT_TO_EMU), int(h_pt * _PT_TO_EMU)


def _bbox_to_emu(
    bbox_x: float,
    bbox_y: float,
    bbox_w: float,
    bbox_h: float,
    slide_w: int,
    slide_h: int,
) -> tuple[int, int, int, int]:
    left = int(slide_w * bbox_x)
    top = int(slide_h * bbox_y)
    width = int(slide_w * max(0.015, bbox_w))
    height = int(slide_h * max(0.012, bbox_h))
    return left, top, width, height


def _line_height_pt(block: TextBlock, slide_h_pt: float) -> float:
    return max(6.0, block.bbox.h * slide_h_pt)


def _normalize_font_sizes(blocks: list[TextBlock], slide_h_pt: float) -> dict[str, float]:
    """One font size per line band — avoids mixed sizes in the same paragraph."""
    sizes: dict[str, float] = {}
    bands: dict[int, list[TextBlock]] = {}
    for block in blocks:
        band = int(block.bbox.y * 200)
        bands.setdefault(band, []).append(block)

    for band_blocks in bands.values():
        heights = [_line_height_pt(b, slide_h_pt) for b in band_blocks]
        median = sorted(heights)[len(heights) // 2]
        cap = 11.0 if any(b.role == "caption" for b in band_blocks) else 18.0
        floor = 7.0 if any(b.role == "caption" for b in band_blocks) else 9.0
        size = max(floor, min(cap, median * 0.88))
        for block in band_blocks:
            sizes[block.id] = size
    return sizes


def _font_size_pt(block: TextBlock, slide_h_pt: float, normalized: dict[str, float]) -> float:
    if block.id in normalized:
        return normalized[block.id]
    if block.font_size_pt and block.font_size_pt > 7:
        cap = 11.0 if block.role == "caption" else 22.0
        return max(7.0, min(cap, block.font_size_pt))
    base = _line_height_pt(block, slide_h_pt) * 0.88
    cap = 11.0 if block.role == "caption" else 22.0
    floor = 7.0 if block.role == "caption" else 9.0
    return max(floor, min(cap, base))


def _center_in(bbox: Bbox, container: Bbox) -> bool:
    cx = bbox.x + bbox.w / 2
    cy = bbox.y + bbox.h / 2
    return (
        container.x <= cx <= container.x + container.w
        and container.y <= cy <= container.y + container.h
    )


def _texts_on_page(ir: DocumentIR, page_index: int) -> list[TextBlock]:
    return [
        b
        for b in ir.blocks
        if isinstance(b, TextBlock) and b.page_index == page_index
    ]


def _image_labels_for_block(
    block: ImageBlock, page_texts: list[TextBlock]
) -> list[TextBlock]:
    labels: list[TextBlock] = []
    for tb in page_texts:
        if _center_in(tb.bbox, block.bbox):
            rel = Bbox(
                x=(tb.bbox.x - block.bbox.x) / max(block.bbox.w, 0.001),
                y=(tb.bbox.y - block.bbox.y) / max(block.bbox.h, 0.001),
                w=tb.bbox.w / max(block.bbox.w, 0.001),
                h=tb.bbox.h / max(block.bbox.h, 0.001),
            )
            labels.append(
                TextBlock(
                    page_index=tb.page_index,
                    bbox=rel,
                    text=tb.text,
                    role=tb.role,
                    font_size_pt=tb.font_size_pt,
                    font_name=tb.font_name,
                )
            )
    return labels


def _blocks_for_page(ir: DocumentIR, page_index: int) -> list:
    blocks = []
    for block in ir.blocks:
        if block.type == "page_break":
            continue
        if isinstance(block, ImageBlock) and _is_gap_path(block.image_path):
            continue
        if getattr(block, "page_index", -1) == page_index:
            blocks.append(block)
    return blocks


def _style_textbox(
    txbox,
    text: str,
    size_pt: float,
    *,
    bold: bool = False,
    caption: bool = False,
) -> None:
    txbox.fill.solid()
    txbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
    txbox.line.fill.background()
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT
    if bold:
        p.font.bold = True
    if caption:
        p.font.size = Pt(min(size_pt, 10.0))


def emit_pptx(
    ir: DocumentIR,
    output_path: Path,
    *,
    page_backgrounds: list[Path] | None = None,
    mask_backgrounds: bool = True,
    mask_figure_labels: bool = True,
) -> Path:
    """Masked scan background + editable text boxes (body + figure captions)."""
    prs = Presentation()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    page_count = ir.source.page_count
    pages = ir.pages or [Page(index=i) for i in range(page_count)]
    if len(pages) < page_count:
        pages = pages + [Page(index=i) for i in range(len(pages), page_count)]

    first = pages[0] if pages else Page(index=0)
    slide_w, slide_h = _page_size(first, 720.0, 405.0)
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    blank_layout = prs.slide_layouts[6]
    mask_dir = output_path.parent / "masked_assets"
    mask_dir.mkdir(parents=True, exist_ok=True)

    for pi in range(page_count):
        slide = prs.slides.add_slide(blank_layout)
        page = pages[pi] if pi < len(pages) else None
        slide_h_pt = page.height_pt if page else 405.0
        sw, sh = slide_w, slide_h

        page_texts = _texts_on_page(ir, pi)
        font_sizes = _normalize_font_sizes(page_texts, slide_h_pt)
        page_blocks = _blocks_for_page(ir, pi)
        images = [b for b in page_blocks if isinstance(b, ImageBlock)]
        tables = [b for b in page_blocks if isinstance(b, TableBlock)]

        bg_path: Path | None = None
        if page and page.background_image:
            candidate = Path(page.background_image)
            if candidate.exists():
                bg_path = candidate
        if not bg_path and page_backgrounds and pi < len(page_backgrounds):
            candidate = page_backgrounds[pi]
            if candidate.exists():
                bg_path = candidate

        if bg_path:
            use_bg = bg_path
            if mask_backgrounds and page_texts:
                masked = mask_dir / f"page_{pi:03d}_masked.png"
                use_bg = mask_page_background(bg_path, page_texts, masked)
            try:
                slide.shapes.add_picture(str(use_bg), 0, 0, width=sw, height=sh)
            except Exception:
                pass

        if not bg_path:
            for block in images:
                if not block.image_path or not Path(block.image_path).exists():
                    continue
                img_path = Path(block.image_path)
                if mask_figure_labels:
                    rel_labels = _image_labels_for_block(block, page_texts)
                    if rel_labels:
                        masked_img = mask_dir / f"{img_path.stem}_masked.png"
                        img_path = mask_text_regions(img_path, rel_labels, masked_img, pad=0.006)

                left, top, width, height = _bbox_to_emu(
                    block.bbox.x, block.bbox.y, block.bbox.w, block.bbox.h, sw, sh
                )
                try:
                    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
                except Exception:
                    pass

        for block in tables:
            rows = block.rows
            if not rows:
                continue
            cols = max(len(r) for r in rows)
            left, top, width, height = _bbox_to_emu(
                block.bbox.x, block.bbox.y, max(block.bbox.w, 0.3), max(block.bbox.h, 0.15), sw, sh
            )
            shape = slide.shapes.add_table(len(rows), cols, left, top, width, height)
            table = shape.table
            for ri, row in enumerate(rows):
                for ci in range(cols):
                    cell_text = row[ci] if ci < len(row) else ""
                    table.cell(ri, ci).text = clean_ocr_text(cell_text)

        body = [b for b in page_texts if b.role != "caption"]
        captions = [b for b in page_texts if b.role == "caption"]

        for block in body + captions:
            text = clean_ocr_text(block.text)
            if not text.strip():
                continue
            left, top, width, height = _bbox_to_emu(
                block.bbox.x, block.bbox.y, block.bbox.w, block.bbox.h, sw, sh
            )
            txbox = slide.shapes.add_textbox(left, top, width, height)
            size = _font_size_pt(block, slide_h_pt, font_sizes)
            _style_textbox(
                txbox,
                text,
                size,
                bold=block.role in ("title", "heading"),
                caption=block.role == "caption",
            )

    prs.save(output_path)
    return output_path
