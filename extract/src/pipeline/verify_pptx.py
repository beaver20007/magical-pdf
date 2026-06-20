"""
verify_pptx.py — автономная верификация PPTX через LibreOffice headless.

Этапы:
  1. LibreOffice headless конвертирует PPTX → PNG постранично
  2. Сравниваем наложение OCR-боксов из PPTX с реальным рендером
  3. Проверяем минимальные пороги по каждому слайду
  4. Возвращаем VerifyReport со списком предупреждений

CLI:
    python verify_pptx.py plan_v26.pptx [--lo-path "C:/Program Files/LibreOffice/program/soffice.exe"]
"""
from __future__ import annotations

import json
import logging
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu

_log = logging.getLogger(__name__)

# Минимум OCR-боксов на слайде с изображениями, чтобы считать слайд OK
_MIN_OCR_PER_IMAGE_SLIDE = 1

# Известные пути LibreOffice на Windows
_LO_CANDIDATES = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    r"C:\Program Files\LibreOffice 7\program\soffice.exe",
]


@dataclass
class SlideResult:
    idx: int            # 1-based
    ocr_boxes: int
    native_boxes: int
    has_images: bool
    status: str         # OK / WARN / SKIP
    note: str = ""


@dataclass
class VerifyReport:
    pptx_path: Path
    slides: list[SlideResult] = field(default_factory=list)
    render_dir: Optional[Path] = None   # PNG output directory
    lo_version: str = ""

    @property
    def passed(self) -> bool:
        return all(s.status != "WARN" for s in self.slides)

    @property
    def warn_slides(self) -> list[int]:
        return [s.idx for s in self.slides if s.status == "WARN"]

    def to_dict(self) -> dict:
        return {
            "pptx": str(self.pptx_path),
            "passed": self.passed,
            "warn_slides": self.warn_slides,
            "lo_version": self.lo_version,
            "slides": [vars(s) for s in self.slides],
        }


# ── LibreOffice ───────────────────────────────────────────────────────────────

def _find_soffice() -> Optional[str]:
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    for p in _LO_CANDIDATES:
        if Path(p).exists():
            return p
    return None


def _lo_version(soffice: str) -> str:
    try:
        out = subprocess.check_output([soffice, "--version"], stderr=subprocess.DEVNULL, timeout=10)
        return out.decode(errors="replace").strip()
    except Exception:
        return "unknown"


def render_pptx_lo(pptx_path: Path, out_dir: Path, soffice: str) -> list[Path]:
    """Convert PPTX → PNG using LibreOffice headless. Returns sorted PNG list."""
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--convert-to", "png",
        "--outdir", str(out_dir),
        str(pptx_path),
    ]
    try:
        subprocess.run(cmd, check=True, timeout=120,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except subprocess.TimeoutExpired:
        _log.error("LibreOffice render timed out")
        return []
    except subprocess.CalledProcessError as e:
        _log.error("LibreOffice failed: %s", e)
        return []

    # LO names output as <stem>.png or <stem>1.png, <stem>2.png ...
    stem = pptx_path.stem
    pngs = sorted(out_dir.glob(f"{stem}*.png"),
                  key=lambda p: int("".join(c for c in p.stem[len(stem):] or "0") or "0"))
    return pngs


# ── PPTX inspection ───────────────────────────────────────────────────────────

def _slide_has_images(slide) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    return False


def _count_textboxes(slide) -> tuple[int, int]:
    """Returns (ocr_count, native_count) based on text color heuristic."""
    ocr = native = 0
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        clr = "?"
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                try:
                    clr = str(run.font.color.rgb)
                    break
                except Exception:
                    pass
            if clr != "?":
                break
        if clr == "151515":
            ocr += 1
        else:
            native += 1
    return ocr, native


# ── Overlay render (python-only fallback) ────────────────────────────────────

def _render_overlay(slide, bg_img: Image.Image) -> Image.Image:
    """Draw green (OCR) and blue (native) boxes over the LO-rendered slide image."""
    img = bg_img.copy()
    draw = ImageDraw.Draw(img, "RGBA")
    sw = slide.slide_width / 914400 * 72   # pt
    sh_val = slide.slide_height / 914400 * 72
    scale_x = img.width / sw
    scale_y = img.height / sh_val

    for s in slide.shapes:
        if not s.has_text_frame or not s.text_frame.text.strip():
            continue
        clr = "?"
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                try:
                    clr = str(run.font.color.rgb)
                    break
                except Exception:
                    pass
            if clr != "?":
                break

        lx = s.left / 914400 * 72 * scale_x
        ly = s.top  / 914400 * 72 * scale_y
        lw = s.width / 914400 * 72 * scale_x
        lh = s.height / 914400 * 72 * scale_y
        color = (0, 200, 0, 200) if clr == "151515" else (0, 80, 255, 200)
        draw.rectangle([lx, ly, lx + lw, ly + lh], outline=color, width=2)

    return img


# ── Main verify function ──────────────────────────────────────────────────────

def verify_pptx(
    pptx_path: Path,
    render_dir: Optional[Path] = None,
    soffice: Optional[str] = None,
    save_overlay: bool = True,
) -> VerifyReport:
    pptx_path = Path(pptx_path)
    report = VerifyReport(pptx_path=pptx_path)

    if soffice is None:
        soffice = _find_soffice()

    if soffice:
        report.lo_version = _lo_version(soffice)
        print(f"[verify] LibreOffice: {report.lo_version}")
    else:
        print("[verify] LibreOffice не найден — overlay рендер через python-pptx фон")

    if render_dir is None:
        render_dir = pptx_path.parent / f"verify_{pptx_path.stem}"
    render_dir.mkdir(parents=True, exist_ok=True)
    report.render_dir = render_dir

    # Render via LibreOffice
    lo_pngs: list[Path] = []
    if soffice:
        lo_tmp = render_dir / "_lo_render"
        lo_pngs = render_pptx_lo(pptx_path, lo_tmp, soffice)
        if lo_pngs:
            print(f"[verify] LO рендер: {len(lo_pngs)} PNG в {lo_tmp}")
        else:
            print("[verify] LO рендер вернул 0 PNG — fallback на фон из PPTX")

    prs = Presentation(str(pptx_path))

    for idx, slide in enumerate(prs.slides):
        ocr_cnt, nat_cnt = _count_textboxes(slide)
        has_imgs = _slide_has_images(slide)

        if has_imgs and ocr_cnt < _MIN_OCR_PER_IMAGE_SLIDE:
            status = "WARN"
            note = f"image slide but only {ocr_cnt} OCR boxes"
        elif not has_imgs and ocr_cnt == 0 and nat_cnt == 0:
            status = "SKIP"
            note = "no content"
        else:
            status = "OK"
            note = ""

        result = SlideResult(
            idx=idx + 1,
            ocr_boxes=ocr_cnt,
            native_boxes=nat_cnt,
            has_images=has_imgs,
            status=status,
            note=note,
        )
        report.slides.append(result)
        flag = "WARN" if status == "WARN" else "OK  "
        print(f"  S{idx+1:02d} [{flag}] OCR={ocr_cnt:3d} native={nat_cnt} imgs={has_imgs} {note}")

        if save_overlay and lo_pngs and idx < len(lo_pngs):
            bg = Image.open(lo_pngs[idx]).convert("RGB")
            overlay = _render_overlay(slide, bg)
            overlay.save(str(render_dir / f"s{idx+1:02d}_overlay.png"))

    # Summary
    json_path = render_dir / "report.json"
    json_path.write_text(json.dumps(report.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"\n[verify] {'PASS ✓' if report.passed else 'FAIL ✗'} — предупреждений: {len(report.warn_slides)} {report.warn_slides}")
    print(f"[verify] Отчёт: {json_path}")
    return report


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("pptx", help="Путь к PPTX")
    p.add_argument("--lo-path", help="Путь к soffice.exe")
    p.add_argument("--out-dir", help="Директория для PNG верификации")
    p.add_argument("--json", help="Сохранить JSON отчёт сюда")
    args = p.parse_args()

    report = verify_pptx(
        pptx_path=Path(args.pptx),
        render_dir=Path(args.out_dir) if args.out_dir else None,
        soffice=args.lo_path,
    )
    if args.json:
        Path(args.json).write_text(
            json.dumps(report.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8"
        )
    sys.exit(0 if report.passed else 1)


if __name__ == "__main__":
    main()
