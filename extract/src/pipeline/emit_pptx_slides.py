"""
emit_pptx_slides.py — конвертация slide-PDF в PPTX с редактируемым текстом.

Архитектура (v5):
  1. Рендерим страницу в PNG
  2. Вырезаем каждый image-блок (чертёж, иллюстрацию) как отдельный кроп
  3. На кропе запускаем EasyOCR → находим текстовые области внутри изображения
  4. В кропе замазываем OCR-области цветом фона (убираем текст из растра)
  5. Вектор-слой страницы (без текста и без image-blocks) — отдельный PNG
  6. В PPTX:
       - Фон слайда (цвет страницы)
       - Вектор-слой (линии, рамки, декор)
       - Каждый кроп на своей позиции (замазанный — только графика без текста)
       - Text boxes из native PDF text
       - Text boxes из OCR поверх каждого кропа (на координатах кропа)
  Результат: весь текст редактируем, включая надписи внутри чертежей.
"""
from __future__ import annotations

from pathlib import Path
from statistics import mode as stat_mode
from typing import NamedTuple

import json

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

try:
    from src.pipeline.ocr_cache import load_ocr_cache as _load_ocr_cache, save_ocr_cache as _save_ocr_cache
    _HAS_OCR_CACHE = True
except Exception:
    _HAS_OCR_CACHE = False

try:
    import cv2 as _cv2
    _HAS_CV2 = True
except ImportError:
    _cv2 = None  # type: ignore[assignment]
    _HAS_CV2 = False

_PT_TO_EMU = 12700

_TESSERACT_CMD = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


def _emu(pt: float) -> int:
    return int(pt * _PT_TO_EMU)


# ── Цвет фона ────────────────────────────────────────────────────────────────

def _dominant_color(
    img: Image.Image, x0: int, y0: int, x1: int, y1: int, band: int = 8
) -> tuple[int, int, int]:
    w, h = img.size
    pixels: list[tuple[int, int, int]] = []
    for px in range(max(0, x0 - band), min(w, x1 + band)):
        for py in [max(0, y0 - band - 1), min(h - 1, y1 + band)]:
            pixels.append(img.getpixel((px, py))[:3])
    for py in range(max(0, y0 - band), min(h, y1 + band)):
        for px in [max(0, x0 - band - 1), min(w - 1, x1 + band)]:
            pixels.append(img.getpixel((px, py))[:3])
    if not pixels:
        return (255, 255, 255)
    bucketed = [(r // 32 * 32, g // 32 * 32, b // 32 * 32) for r, g, b in pixels]
    try:
        dom = stat_mode(bucketed)
    except Exception:
        dom = bucketed[0]
    in_bkt = [p for p in pixels if (p[0]//32*32, p[1]//32*32, p[2]//32*32) == dom]
    return (
        sum(p[0] for p in in_bkt) // len(in_bkt),
        sum(p[1] for p in in_bkt) // len(in_bkt),
        sum(p[2] for p in in_bkt) // len(in_bkt),
    )


def _page_bg_color(img: Image.Image) -> tuple[int, int, int]:
    w, h = img.size
    corners = [(2, 2), (w-3, 2), (2, h-3), (w-3, h-3)]
    colors = [img.getpixel(c)[:3] for c in corners]
    return (
        sum(c[0] for c in colors) // 4,
        sum(c[1] for c in colors) // 4,
        sum(c[2] for c in colors) // 4,
    )


def _mask_rects(
    img: Image.Image,
    rects_px: list[tuple[int, int, int, int]],
    fill: tuple[int, int, int] | None = None,
    pad: int = 2,
) -> Image.Image:
    """Замазываем прямоугольники (px) цветом фона или заданным цветом."""
    img = img.copy()
    draw = ImageDraw.Draw(img)
    iw, ih = img.size
    for x0, y0, x1, y1 in rects_px:
        px0 = max(0, x0 - pad)
        py0 = max(0, y0 - pad)
        px1 = min(iw, x1 + pad)
        py1 = min(ih, y1 + pad)
        if px1 <= px0 or py1 <= py0:
            continue
        c = fill if fill else _dominant_color(img, px0, py0, px1, py1)
        draw.rectangle([px0, py0, px1, py1], fill=c)
    return img


def _pt_rects_to_px(
    rects_pt: list[tuple], scale: float
) -> list[tuple[int, int, int, int]]:
    return [
        (int(x0*scale), int(y0*scale), int(x1*scale), int(y1*scale))
        for x0, y0, x1, y1 in rects_pt
    ]


# ── OCR на кропе ─────────────────────────────────────────────────────────────

class OcrWord(NamedTuple):
    text: str
    # Координаты в пикселях КРОПА
    px0: int
    py0: int
    px1: int
    py1: int
    conf: float


import re as _re
_ALPHANUM      = _re.compile(r"[А-Яа-яёЁA-Za-z0-9]")
_CYRILLIC      = _re.compile(r"[А-Яа-яёЁ]")
# "чистые" символы — всё что может быть в реальной аннотации чертежа
_CLEAN_ALNUM   = _re.compile(r"[А-Яа-яёЁA-Za-z0-9\-—.,:%°+=()①-⑨₁-₉/\\]")
_NOISE         = _re.compile(r"^[|\\/*+=(){}\[\]<>«»\"\'`^~@#$%&]{1,2}$")
# Паттерны измерительных значений: —0.36, -1.13, ±0.05, .30 (ведущая точка), =0.36
_MEASUREMENT   = _re.compile(r"^([—\-±=]?\d+[.,]\d+|[.,]\d+)$")
# Phase 61: 2-decimal float tokens subject to Float Re-read
_FLOAT_2DEC    = _re.compile(r'^-?\d+[.,]\d{2}$')
# Технические обозначения чертежей: i=4%, 1:200, 0.5%, уклон-значения
_TECH_MARKER   = _re.compile(r"^(i=[0-9]+%?|[0-9]+:[0-9]+|[0-9]+[.,][0-9]+%|[А-Яа-яёЁ]{2,10}[0-9]*)$")
# Штриховка читается как короткие латинские "слова" из e,s,a,f,t,i,c (без кириллицы/цифр)
# Строчные латинские слова ≤5 символов — штриховка (ah, ar, ee, jam, tum, bema).
# Заглавные (Dn, PP) и смешанные не попадают под этот фильтр.
_SHORT_LOWER_LAT = _re.compile(r"^[a-z]{2,5}[).,!|]*$")
# Полностью заглавные латинские слова >=3 букв — OCR мусор (ALN, ALY, GRE, RAS, ALLY…)
# 2-буквенные (PP, DN, SS, CA) обрабатываются отдельно в inline-фильтре
_ALL_CAPS_LATIN  = _re.compile(r"^[A-Z]{3,}$")
_DASH_SEQ        = _re.compile(r'^[\-—<>/\\]{2,6}$')
_ROT_MIN_CONF    = 55
_TECH_WHITELIST = frozenset({
    'мм', 'см', 'дм', 'км', 'кг', 'шт', 'пм', 'пл', 'вп', 'ту',
    'пвх', 'пп', 'дп', 'бв', 'пу', 'ду', 'дн', 'кн', 'мн', 'па',
    'dn', 'pp', 'pe', 'kg', 'mm', 'cm', 'wt', 'kw', 'hz', 'pvc',
    'pp-r', 'pe-x', 'pvc-u', 'pvc-c', 'hdpe',
})
# Смешанные буквенно-цифровые технические токены: Dn100, DN50, PN16, PVC32, Dn110
# Цифры обязательны — чтобы не захватывать чисто-буквенный мусор (ALY, MBX, ALLY)
_TECH_MIXED = _re.compile(r'^[A-Za-z]{1,4}[-]?[A-Za-z]{0,2}\d{2,4}[a-zA-Z]?$')
# Известные бренды/продукты сантехнической и строительной области
_DOMAIN_BRANDS = frozenset({
    'betomax', 'polymax', 'perfokora', 'perfokor', 'wafix', 'wavin',
    'rehau', 'kan', 'gebo', 'viega', 'valtek', 'kalde',
})


def _apply_char_subs(text: str) -> str:
    """
    Исправляет кириллические буквы, ошибочно распознанные в числовом контексте.
    О/о → 0 рядом с цифрами/точкой, З → 3 рядом с цифрами.
    """
    t = text
    # Кириллическая О/о между цифрами или знаками пунктуации (.,)
    t = _re.sub(r'(?<=[.,\d])[Оо](?=[.,\d])', '0', t)
    # Кириллическая О/о в начале числа
    t = _re.sub(r'^[Оо](?=[.,\d])', '0', t)
    # Кириллическая О/о в конце числа
    t = _re.sub(r'(?<=[.,\d])[Оо]$', '0', t)
    # Кириллическая З между цифрами
    t = _re.sub(r'(?<=[\d])З(?=[\d])', '3', t)
    # Кириллическая З в начале числа
    t = _re.sub(r'^З(?=[.,\d])', '3', t)
    return t


_CYRILLIC_ENGINEERING = frozenset({
    'отверстие', 'отверстия', 'фундамент', 'фундаменте', 'ограждение', 'ограждения',
    'трубопровод', 'дренаж', 'колодец', 'люк', 'решетка', 'уклон', 'отметка',
    'существующая', 'проектная', 'граница', 'скважина', 'насос', 'септик',
    'канализация', 'газопровод', 'горизонталь',
    # Phase 66: Cyrillic Engineering Vocabulary Boost
    'водопровод', 'водоснабжение', 'водоотведение', 'теплоснабжение',
    'электроснабжение', 'вентиляция', 'кровля', 'перекрытие', 'перегородка',
    'фасад', 'фундаментная', 'подвал', 'цоколь', 'парапет', 'пандус',
    'лестница', 'пролёт', 'арматура', 'бетон', 'железобетон',
    'монолит', 'сваи', 'ростверк', 'плита', 'балка', 'колонна',
    'стена', 'проём', 'перемычка', 'ниша', 'штраба',
})

_CYRILLIC_SUFFIXES = _re.compile(
    r'(ия|ание|ение|ость|тель|ник|ный|ного|ого)$', _re.IGNORECASE
)


def _validate_cyrillic_token(text: str) -> bool:
    """Проверяет, является ли кириллическое слово валидным инженерным термином.

    Возвращает True если слово:
    - совпадает с известным инженерным термином, ИЛИ
    - длиной >= 5 чисто кириллических символов с известным суффиксом.
    """
    t = text.strip().lower()
    if t in _CYRILLIC_ENGINEERING:
        return True
    if len(t) >= 5 and _CYRILLIC.fullmatch(t):
        if _CYRILLIC_SUFFIXES.search(t):
            return True
    return False


def _is_valid_ocr(text: str) -> bool:
    """
    Принимаем OCR-результат если:
    - длина >= 1 (одиночный маркер типа B, C, А)
    - содержит хотя бы 1 кириллическую/латинскую букву или цифру
    - не является пунктуационным мусором или штриховочным шумом
    - ≥ 40% символов — "чистые"
    """
    t = text.strip()
    if t.lower() in _TECH_WHITELIST:
        return True
    if _TECH_MIXED.match(t) and not _re.search(r'[а-яёА-ЯЁ]', t) and len(t) >= 2:
        return True
    if not t:
        return False
    # Одиночный символ — только размерный маркер (заглавная буква или цифра)
    if len(t) == 1:
        return bool(_re.match(r"[А-ЯA-Z0-9]", t))
    if not _ALPHANUM.search(t):
        return False
    if _NOISE.match(t):
        return False
    if _DASH_SEQ.match(t) and not _ALPHANUM.search(t):
        return False
    # Строчные латинские слова ≤5 символов без цифр/кириллицы — штриховка или мусор
    if _SHORT_LOWER_LAT.match(t) and not _CYRILLIC.search(t) and not _re.search(r"\d", t):
        return False
    # Длинные (>8 символов) латинские слова без кириллицы и цифр — OCR мусор
    # (tlecxoynoeumens, Tlecxoynoeumens) — не могут быть реальными словами в рус. чертеже
    if len(t) > 10 and t.isalpha() and not _CYRILLIC.search(t):
        return False
    # Слова оканчивающиеся на }]| — рамки таблиц (ints}, Lda])
    if t[-1] in ('}', ']', '|') and not _CYRILLIC.search(t) and len(t) <= 8:
        return False
    # Полностью заглавные латинские слова ≥4 букв без кириллицы — OCR мусор (NALLY, ALLY)
    if _ALL_CAPS_LATIN.match(t) and not _CYRILLIC.search(t):
        return False
    # Известные бренды/продукты — всегда принимаем независимо от регистра
    if t.lower() in _DOMAIN_BRANDS:
        return True
    # Фильтруем акцентный мусор (ées, éée): доля "чистых" символов ≥ 40%
    no_space = t.replace(" ", "")
    if no_space:
        clean_count = sum(1 for c in no_space if _CLEAN_ALNUM.match(c))
        if clean_count / len(no_space) < 0.40:
            return False
    return True


def _min_conf_for(text: str) -> int:
    """Минимальный confidence в зависимости от типа текста."""
    t = text.strip()
    if text.strip().lower() in _TECH_WHITELIST:
        return 30
    if _TECH_MIXED.match(t.strip()):
        return 30
    if len(t) == 1:
        return 80   # одиночные символы — только при высокой уверенности
    if _MEASUREMENT.match(t):
        return 20   # числа-измерения — очень низкий порог, сложный фон на чертежах
    if _TECH_MARKER.match(t):
        return 20   # технические маркеры: i=4%, 1:200, уклон
    if _re.match(r"^\d{1,4}$", t):
        return 25   # целые числа (размеры, диаметры) — чуть выше порог от мусора
    if len(t) <= 2:
        if _re.match(r'^[а-яёa-z]{2}$', t):
            return 70   # 2-char строчные фрагменты (ух, ен, oe) — повышенный порог
        return 50   # короткие аббревиатуры (заглавные, смешанные)
    if len(t) >= 5 and _CYRILLIC.search(t) and not _re.search(r'[A-Za-z]', t):
        if _validate_cyrillic_token(t):
            # Phase 66: lower threshold for longer well-formed Cyrillic engineering terms
            if len(t) >= 8 and _CYRILLIC.fullmatch(t) and _CYRILLIC_SUFFIXES.search(t):
                return 18
            return 18
        return 35   # длинное кириллическое слово — умеренный порог
    return 28       # всё остальное


def _remove_hatching(gray_arr) -> object:
    """
    Морфологическое opening (PIL) для удаления штриховки.
    Штриховка = тонкие линии (~1-3px); текст = более толстые штрихи (~5-15px).
    Opening (эрозия→дилатация) убирает тонкие объекты, сохраняет толстые.
    """
    import numpy as np
    from PIL import Image as _PILImg, ImageFilter

    pil = _PILImg.fromarray(gray_arr.astype("uint8"))
    # Инвертируем: чёрный текст → белый (для min/max filter)
    inv = _PILImg.fromarray(255 - np.array(pil))
    # Opening: MinFilter (эрозия) затем MaxFilter (дилатация)
    kernel = 5  # размер ядра > толщина штриховки
    eroded  = inv.filter(ImageFilter.MinFilter(kernel))
    dilated = eroded.filter(ImageFilter.MaxFilter(kernel))
    # Инвертируем обратно: белый фон, чёрный текст
    result = _PILImg.fromarray(255 - np.array(dilated))
    return np.array(result)


def _make_gray_variants(crop_img: Image.Image) -> list[tuple[Image.Image, float]]:
    """
    Возвращает несколько вариантов препроцессинга:
    1. Luminosity-grayscale + контраст + 2x upscale
    2. Min(R,G,B) + контраст + 2x upscale — цветной текст (синий, красный) становится чёрным
    Всегда 2x upscale чтобы 4-6pt аннотации (≈20px native) стали ≥40px для Tesseract.
    """
    from PIL import ImageEnhance
    import numpy as _np

    w, h = crop_img.size
    # Для больших изображений уменьшаем масштаб чтобы не съесть всю RAM
    # Tesseract работает хорошо при 150-300 DPI; нативные PDF-изображения уже ~150 DPI
    MAX_DIM = 3000
    if max(w, h) * 2 > MAX_DIM:
        UPSCALE = max(1, MAX_DIM // max(w, h))
    else:
        UPSCALE = 2
    if UPSCALE < 1:
        UPSCALE = 1
    results: list[tuple[Image.Image, float]] = []

    def _up(img: Image.Image) -> Image.Image:
        if UPSCALE == 1:
            return img
        return img.resize((w * UPSCALE, h * UPSCALE), Image.LANCZOS)

    # Вариант 1: стандартный grayscale
    g1 = ImageEnhance.Contrast(crop_img.convert("L")).enhance(2.5)
    results.append((_up(g1).convert("RGB"), float(UPSCALE)))

    # Вариант 2: min(R,G,B) — любой насыщенный цвет → чёрный
    arr = _np.array(crop_img.convert("RGB"))
    min_ch = _np.min(arr, axis=2).astype(_np.uint8)
    g2 = ImageEnhance.Contrast(Image.fromarray(min_ch)).enhance(2.0)
    results.append((_up(g2).convert("RGB"), float(UPSCALE)))

    # Вариант 3: выравнивание гистограммы (ImageOps.equalize) — усиливает низкоконтрастный текст
    from PIL import ImageOps
    g3 = ImageOps.equalize(crop_img.convert("L"))
    g3 = ImageEnhance.Contrast(g3).enhance(1.8)
    results.append((_up(g3).convert("RGB"), float(UPSCALE)))

    # Вариант 4: для плотной штриховки — сильное морфологическое opening (kernel=9)
    # Определяем наличие плотной штриховки по дисперсии пикселей (высокая плотность градиентов).
    gray_arr4 = _np.array(crop_img.convert("L"))
    _variance = float(_np.var(gray_arr4.astype("float32")))
    # Нормализуем: максимальная теоретическая дисперсия = 128^2 = 16384
    _edge_density = _variance / 16384.0
    if _edge_density > 0.30:
        # Плотная штриховка обнаружена: применяем opening с kernel=9
        from PIL import Image as _PILImg, ImageFilter
        _pil4 = _PILImg.fromarray(gray_arr4.astype("uint8"))
        _inv4 = _PILImg.fromarray(255 - _np.array(_pil4))
        _k9 = 9
        _eroded4 = _inv4.filter(ImageFilter.MinFilter(_k9))
        _dilated4 = _eroded4.filter(ImageFilter.MaxFilter(_k9))
        _opened4 = _PILImg.fromarray(255 - _np.array(_dilated4))
        g4 = ImageEnhance.Contrast(_opened4).enhance(2.0)
        results.append((_up(g4).convert("RGB"), float(UPSCALE)))

    # Вариант 5: локальная адаптивная бинаризация — для сканов с неравномерной подсветкой
    # (тени от сгиба, виньетирование). Глобальный контраст не справляется с локально тёмными зонами.
    # Алгоритм: local_mean = BoxBlur(r=16); pixel < local_mean*0.88 → 0 (чёрный), иначе 255 (белый).
    try:
        import numpy as _np5
        from PIL import ImageFilter as _IF5
        _g5 = crop_img.convert("L")
        _local_mean5 = _g5.filter(_IF5.BoxBlur(16))
        _arr5 = _np5.array(_g5, dtype=_np5.float32)
        _mean5 = _np5.array(_local_mean5, dtype=_np5.float32)
        _bin5 = _np5.where(_arr5 < _mean5 * 0.88, 0, 255).astype(_np5.uint8)
        g5 = Image.fromarray(_bin5)
    except Exception:
        from PIL import ImageFilter as _IF5
        _g5 = crop_img.convert("L")
        _local_mean5 = _g5.filter(_IF5.BoxBlur(16))
        _px5 = list(_g5.getdata())
        _mx5 = list(_local_mean5.getdata())
        _bin5 = bytes([0 if p < m * 0.88 else 255 for p, m in zip(_px5, _mx5)])
        g5 = Image.frombytes("L", _g5.size, _bin5)
    results.append((_up(g5).convert("RGB"), float(UPSCALE)))

    return results


def _run_tesseract(img: Image.Image, psm: int, up_scale: float) -> list[OcrWord]:
    """Запускает Tesseract с заданным psm, возвращает слова в исходных пикселях."""
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
    cfg = f"--oem 1 --psm {psm} -l rus+eng"
    try:
        data = pytesseract.image_to_data(img, config=cfg, output_type=pytesseract.Output.DICT)
    except Exception as e:
        print(f"    Tesseract psm={psm} error: {e}")
        return []

    iw, ih = img.size  # размер апскейлнутого изображения
    # Максимальный допустимый размер bbox слова (в нативных пикселях).
    # PSM 6 убран, но оставляем разумный лимит на случай аномальных bbox:
    # слово шире 55% или выше 30% изображения — скорее всего артефакт
    max_word_w = iw * 0.55 / up_scale
    max_word_h = ih * 0.30 / up_scale

    words: list[OcrWord] = []
    for j in range(len(data["text"])):
        txt = _apply_char_subs((data["text"][j] or "").strip())
        conf = int(data["conf"][j])
        if not txt or not _is_valid_ocr(txt):
            continue
        if conf < _min_conf_for(txt):
            continue
        x0 = int(data["left"][j] / up_scale)
        y0 = int(data["top"][j] / up_scale)
        x1 = int((data["left"][j] + data["width"][j]) / up_scale)
        y1 = int((data["top"][j] + data["height"][j]) / up_scale)
        # Отбрасываем bbox слишком большого размера (шум без PSM 6)
        if (x1 - x0) > max_word_w or (y1 - y0) > max_word_h:
            continue
        words.append(OcrWord(text=txt, px0=x0, py0=y0, px1=x1, py1=y1, conf=conf / 100.0))
    return words


def _merge_word_lists(base: list[OcrWord], extra: list[OcrWord], tol: int = 8) -> list[OcrWord]:
    """Добавляет слова из extra, которых нет в base (по центру bbox с допуском tol px)."""
    def _center(w: OcrWord):
        return ((w.px0 + w.px1) / 2, (w.py0 + w.py1) / 2)
    existing = [_center(w) for w in base]
    merged = list(base)
    for w in extra:
        cx, cy = _center(w)
        if not any(abs(cx - ex) < tol and abs(cy - ey) < tol for ex, ey in existing):
            merged.append(w)
            existing.append((cx, cy))
    return merged


def _vote_ocr_results(variant_results: list[list[OcrWord]]) -> list[OcrWord]:
    """Phase 56: Confidence-Weighted Multi-Variant Voting.

    Groups candidate words from all preprocessing variants by approximate
    position (within 10px x, 8px y), scores each candidate text by summing
    confidence values across variants that produced it plus domain bonuses,
    and returns the highest-scoring word per position group.

    Bonuses:
      +0.20  text in _TECH_WHITELIST
      +0.15  text matches _TECH_MIXED
      +0.10  text matches _MEASUREMENT
    """
    if not variant_results:
        return []

    TOL_X = 10
    TOL_Y = 8

    all_words: list[OcrWord] = []
    for words in variant_results:
        all_words.extend(words)

    if not all_words:
        return []

    def _center(w: OcrWord) -> tuple[float, float]:
        return ((w.px0 + w.px1) / 2.0, (w.py0 + w.py1) / 2.0)

    def _bonus(text: str) -> float:
        t = text.strip()
        score = 0.0
        if t.lower() in _TECH_WHITELIST:
            score += 0.20
        if _TECH_MIXED.match(t):
            score += 0.15
        if _MEASUREMENT.match(t):
            score += 0.10
        return score

    groups: list[list[OcrWord]] = []
    group_centers: list[tuple[float, float]] = []

    for w in all_words:
        cx, cy = _center(w)
        matched = -1
        for gi, (gcx, gcy) in enumerate(group_centers):
            if abs(cx - gcx) <= TOL_X and abs(cy - gcy) <= TOL_Y:
                matched = gi
                break
        if matched >= 0:
            groups[matched].append(w)
            n = len(groups[matched])
            old_gcx, old_gcy = group_centers[matched]
            group_centers[matched] = (
                (old_gcx * (n - 1) + cx) / n,
                (old_gcy * (n - 1) + cy) / n,
            )
        else:
            groups.append([w])
            group_centers.append((cx, cy))

    winners: list[OcrWord] = []
    for group in groups:
        scores: dict[str, float] = {}
        best_word_for: dict[str, OcrWord] = {}
        for w in group:
            t = w.text.strip()
            if not t:
                continue
            s = scores.get(t, 0.0) + w.conf + _bonus(t)
            scores[t] = s
            if t not in best_word_for or w.conf > best_word_for[t].conf:
                best_word_for[t] = w
        if not scores:
            continue
        best_text = max(scores, key=lambda k: scores[k])
        same = [w for w in group if w.text.strip() == best_text]
        avg_px0 = int(sum(w.px0 for w in same) / len(same))
        avg_py0 = int(sum(w.py0 for w in same) / len(same))
        avg_px1 = int(sum(w.px1 for w in same) / len(same))
        avg_py1 = int(sum(w.py1 for w in same) / len(same))
        avg_conf = scores[best_text] / max(len(same), 1)
        avg_conf = min(avg_conf, 1.0)
        winners.append(OcrWord(
            text=best_text,
            px0=avg_px0, py0=avg_py0,
            px1=avg_px1, py1=avg_py1,
            conf=avg_conf,
        ))

    return winners


def _has_dense_text_grid(img: Image.Image) -> bool:
    """
    Определяет плотную текстовую сетку (таблицы, штамп-блоки).
    Два критерия:
    1. >18% тёмных пикселей (значение < 128) от общего числа пикселей.
    2. >= 5 строк пикселей, где >40% пикселей тёмные.
    """
    gray = img.convert("L")
    try:
        import numpy as _np
        arr = _np.array(gray)
        dark_mask = arr < 128
        total = dark_mask.size
        dark_count = int(dark_mask.sum())
        # Критерий 1: общая плотность тёмных пикселей
        if total > 0 and dark_count / total > 0.18:
            return True
        # Критерий 2: количество строк с >40% тёмных пикселей
        row_dark = dark_mask.sum(axis=1)
        row_total = arr.shape[1]
        dense_rows = int((row_dark > row_total * 0.40).sum())
        if dense_rows >= 5:
            return True
    except ImportError:
        # Fallback без numpy
        pixels = list(gray.getdata())
        total = len(pixels)
        if total == 0:
            return False
        dark_count = sum(1 for p in pixels if p < 128)
        if dark_count / total > 0.18:
            return True
        w, h = gray.size
        dense_rows = 0
        for row in range(h):
            row_pixels = pixels[row * w:(row + 1) * w]
            row_dark = sum(1 for p in row_pixels if p < 128)
            if w > 0 and row_dark / w > 0.40:
                dense_rows += 1
        if dense_rows >= 5:
            return True
    return False


def _find_small_regions(words: list[OcrWord]) -> list[tuple[int, int, int, int]]:
    """Находит кластеры маленьких низкоуверенных слов для 4x upscale re-OCR."""
    small = [w for w in words if (w.py1 - w.py0) < 20 and w.conf < 0.60]
    if not small:
        return []
    # Сортируем по px0
    small = sorted(small, key=lambda w: w.px0)
    clusters: list[list[OcrWord]] = []
    cur: list[OcrWord] = [small[0]]
    for w in small[1:]:
        prev = cur[-1]
        prev_cy = (prev.py0 + prev.py1) / 2
        cur_cy  = (w.py0 + w.py1) / 2
        if (w.px0 - prev.px1) < 50 and abs(cur_cy - prev_cy) <= 30:
            cur.append(w)
        else:
            clusters.append(cur)
            cur = [w]
    clusters.append(cur)
    result: list[tuple[int, int, int, int]] = []
    for cl in clusters:
        x0 = min(w.px0 for w in cl) - 8
        y0 = min(w.py0 for w in cl) - 8
        x1 = max(w.px1 for w in cl) + 8
        y1 = max(w.py1 for w in cl) + 8
        result.append((x0, y0, x1, y1))
    return result


def _try_rotated_ocr(crop_img: Image.Image, up_scale: float) -> list[OcrWord]:
    """Пробует OCR с поворотом на 90 и 270 градусов для вертикального текста.
    Возвращает слова с координатами, трансформированными обратно в исходное пространство кропа.
    """
    from PIL import ImageEnhance
    result: list[OcrWord] = []
    for angle in [90, 270]:
        rotated = crop_img.rotate(angle, expand=True)
        g = ImageEnhance.Contrast(rotated.convert("L")).enhance(2.5)
        g_rgb = g.convert("RGB")
        words = _run_tesseract(g_rgb, psm=11, up_scale=up_scale)
        rot_words = [w for w in words if (w.conf if hasattr(w, 'conf') else 50) >= _ROT_MIN_CONF]
        words = rot_words
        if len(words) < 2:
            continue
        avg_conf = sum(w.conf for w in words) / len(words)
        if avg_conf <= 0.55:
            continue
        rw, rh = rotated.size  # после expand=True: rw=orig_h, rh=orig_w (для 90 и 270)
        transformed: list[OcrWord] = []
        for w in words:
            px0, py0, px1, py1 = w.px0, w.py0, w.px1, w.py1
            if angle == 90:
                # rotate(90, expand=True) поворачивает CCW:
                # rotated.size = (orig_h, orig_w)
                # новые координаты в orig: x' = py0, y' = rh - px1, x1' = py1, y1' = rh - px0
                new_px0 = py0
                new_py0 = rh - px1
                new_px1 = py1
                new_py1 = rh - px0
            else:  # 270
                # rotate(270, expand=True) поворачивает CW:
                # rotated.size = (orig_h, orig_w)
                # новые координаты в orig: x' = rw - py1, y' = px0, x1' = rw - py0, y1' = px1
                new_px0 = rw - py1
                new_py0 = px0
                new_px1 = rw - py0
                new_py1 = px1
            transformed.append(OcrWord(
                text=w.text,
                px0=new_px0, py0=new_py0,
                px1=new_px1, py1=new_py1,
                conf=w.conf,
            ))
        result = _merge_word_lists(result, transformed, tol=20)
    return result



def _reread_low_conf_measurements(words: list[OcrWord], crop_img: Image.Image) -> list[OcrWord]:
    """Phase 43: 150% zoom targeted re-read for low-confidence measurement tokens.

    Finds words matching _MEASUREMENT with conf < 0.55 and re-OCRs them at 3x
    zoom using PSM 8 (single word). Replaces the original word if confidence
    improves by more than 0.10. At most 8 re-reads per image to bound runtime.
    """
    iw, ih = crop_img.size
    replaced: dict[int, OcrWord] = {}  # index -> replacement
    re_reads = 0
    for idx, w in enumerate(words):
        if re_reads >= 8:
            break
        if not _MEASUREMENT.match(w.text.strip()):
            continue
        if w.conf >= 0.55:
            continue
        # Crop with 4px padding, clamped to image bounds
        x0 = max(0, w.px0 - 4)
        y0 = max(0, w.py0 - 4)
        x1 = min(iw, w.px1 + 4)
        y1 = min(ih, w.py1 + 4)
        if x1 <= x0 or y1 <= y0:
            continue
        tile = crop_img.crop((x0, y0, x1, y1))
        # 3x upscale -- main pass already does upscaling; this is a targeted re-read
        tile3x = tile.resize((tile.width * 3, tile.height * 3), Image.LANCZOS).convert("L")
        new_words = _run_tesseract(tile3x.convert("RGB"), psm=8, up_scale=3.0)
        if not new_words:
            re_reads += 1
            continue
        # Pick the highest-confidence result
        best = max(new_words, key=lambda nw: nw.conf)
        if best.conf > w.conf + 0.10:
            # Coordinates remapped back to crop_img space (original bbox preserved)
            replaced[idx] = OcrWord(
                best.text,
                w.px0, w.py0, w.px1, w.py1,
                best.conf,
            )
        re_reads += 1

    if not replaced:
        return words
    return [replaced.get(i, w) for i, w in enumerate(words)]


# ── Phase 61: Float Re-read for uncertain decimal values ─────────────────────

def _reread_uncertain_floats(words: list[OcrWord], img: Image.Image) -> list[OcrWord]:
    """Phase 61: High-zoom PSM-8 re-read for 2-decimal float tokens.

    Finds words matching _FLOAT_2DEC (e.g. 0.95, -1.23, 0,88) regardless of
    confidence and re-OCRs them at 1.5× zoom with PSM 8 and a digit-only
    whitelist.  If the new read differs from the original it replaces the
    token (preserving all other OcrWord fields).  At most 12 re-reads per
    image to bound runtime.
    """
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD

    out: list[OcrWord] = []
    re_reads = 0
    iw, ih = img.size
    for w in words:
        if re_reads >= 12 or not _FLOAT_2DEC.match(w.text.strip()):
            out.append(w)
            continue
        try:
            pad = 4
            x0 = max(0, w.px0 - pad)
            y0 = max(0, w.py0 - pad)
            x1 = min(iw, w.px1 + pad)
            y1 = min(ih, w.py1 + pad)
            if x1 <= x0 or y1 <= y0:
                out.append(w)
                continue
            tile = img.crop((x0, y0, x1, y1))
            nw = int(tile.width * 1.5)
            nh = int(tile.height * 1.5)
            if nw < 4 or nh < 4:
                out.append(w)
                continue
            tile = tile.resize((nw, nh), Image.LANCZOS).convert('L').convert('RGB')
            cfg = '--oem 1 --psm 8 --tessedit_char_whitelist 0123456789.,-'
            raw = pytesseract.image_to_string(tile, config=cfg, lang='rus+eng').strip()
            # normalise decimal separator
            raw = raw.replace(',', '.')
            # strip trailing punctuation that sometimes leaks through
            raw = raw.rstrip('.,;')
            if _FLOAT_2DEC.match(raw) and raw != w.text.strip():
                w = w._replace(text=raw)
        except Exception:
            pass
        re_reads += 1
        out.append(w)
    return out


# ── Phase 55: Table Structure Recovery ───────────────────────────────────────

def _detect_table_cells(img: Image.Image) -> list[tuple[int, int, int, int]] | None:
    """Detects table cell bounding boxes via Hough line detection (cv2).

    Converts to grayscale, applies adaptive threshold, runs HoughLinesP to
    find horizontal and vertical line segments.  Returns a list of
    (x0, y0, x1, y1) cell rectangles built from the intersection grid when
    >= 4 horizontal lines AND >= 2 vertical lines are found, otherwise None.
    Returns None immediately when cv2 is not available.

    Conservative Hough params:
      minLineLength(H) = img.width  * 0.3
      minLineLength(V) = img.height * 0.2
      maxLineGap       = 10
      threshold        = 80
    """
    if not _HAS_CV2:
        return None

    import numpy as _np

    iw, ih = img.size
    gray = _np.array(img.convert("L"))

    thresh = _cv2.adaptiveThreshold(
        gray, 255,
        _cv2.ADAPTIVE_THRESH_MEAN_C,
        _cv2.THRESH_BINARY_INV,
        blockSize=15, C=10,
    )

    min_h_len = max(10, int(iw * 0.3))
    min_v_len = max(10, int(ih * 0.2))
    hough_thresh = 80
    max_gap = 10

    h_kernel = _cv2.getStructuringElement(_cv2.MORPH_RECT, (min_h_len // 2, 1))
    h_img = _cv2.morphologyEx(thresh, _cv2.MORPH_OPEN, h_kernel)
    h_lines_raw = _cv2.HoughLinesP(
        h_img, rho=1, theta=3.14159265358979 / 180,
        threshold=hough_thresh,
        minLineLength=min_h_len,
        maxLineGap=max_gap,
    )

    v_kernel = _cv2.getStructuringElement(_cv2.MORPH_RECT, (1, min_v_len // 2))
    v_img = _cv2.morphologyEx(thresh, _cv2.MORPH_OPEN, v_kernel)
    v_lines_raw = _cv2.HoughLinesP(
        v_img, rho=1, theta=3.14159265358979 / 180,
        threshold=hough_thresh,
        minLineLength=min_v_len,
        maxLineGap=max_gap,
    )

    h_ys: list[int] = []
    if h_lines_raw is not None:
        for seg in h_lines_raw:
            x1r, y1r, x2r, y2r = seg[0]
            if abs(y2r - y1r) <= 5:
                h_ys.append((y1r + y2r) // 2)

    v_xs: list[int] = []
    if v_lines_raw is not None:
        for seg in v_lines_raw:
            x1r, y1r, x2r, y2r = seg[0]
            if abs(x2r - x1r) <= 5:
                v_xs.append((x1r + x2r) // 2)

    if len(h_ys) < 4 or len(v_xs) < 2:
        return None

    def _cluster(vals: list[int], gap: int = 8) -> list[int]:
        vals = sorted(set(vals))
        clusters: list[list[int]] = []
        cur = [vals[0]]
        for v in vals[1:]:
            if v - cur[-1] <= gap:
                cur.append(v)
            else:
                clusters.append(cur)
                cur = [v]
        clusters.append(cur)
        return [sum(c) // len(c) for c in clusters]

    h_sorted = _cluster(h_ys)
    v_sorted = _cluster(v_xs)

    if len(h_sorted) < 4 or len(v_sorted) < 2:
        return None

    cells: list[tuple[int, int, int, int]] = []
    for ri in range(len(h_sorted) - 1):
        for ci in range(len(v_sorted) - 1):
            x0 = v_sorted[ci]
            y0 = h_sorted[ri]
            x1 = v_sorted[ci + 1]
            y1 = h_sorted[ri + 1]
            if x1 - x0 < 4 or y1 - y0 < 4:
                continue
            cells.append((x0, y0, x1, y1))

    return cells if cells else None


def _ocr_table_cells(
    crop_img: Image.Image,
    cells: list[tuple[int, int, int, int]],
) -> list[OcrWord]:
    """OCR each table cell individually with PSM 7 (single text line).

    Returns OcrWords with pixel coordinates mapped back to the full crop image
    space (not the individual cell space).
    """
    from PIL import ImageEnhance as _IE

    words: list[OcrWord] = []
    iw, ih = crop_img.size

    for x0, y0, x1, y1 in cells:
        cx0 = max(0, x0 - 2)
        cy0 = max(0, y0 - 2)
        cx1 = min(iw, x1 + 2)
        cy1 = min(ih, y1 + 2)
        if cx1 <= cx0 or cy1 <= cy0:
            continue
        cell_img = crop_img.crop((cx0, cy0, cx1, cy1))
        cw, ch = cell_img.size
        if cw < 4 or ch < 4:
            continue
        up = 3
        cell_up = cell_img.resize((cw * up, ch * up), Image.LANCZOS)
        cell_up = _IE.Contrast(cell_up.convert("L")).enhance(2.5).convert("RGB")
        raw_words = _run_tesseract(cell_up, psm=7, up_scale=float(up))
        for w in raw_words:
            words.append(OcrWord(
                w.text,
                cx0 + w.px0,
                cy0 + w.py0,
                cx0 + w.px1,
                cy0 + w.py1,
                w.conf,
            ))

    return words


def _ocr_crop(crop_img: Image.Image, assets_dir=None) -> list[OcrWord]:
    """
    OCR через Tesseract v5 (LSTM).
    Запускает несколько вариантов препроцессинга (grayscale + min-channel)
    и режимов PSM (11=sparse, 6=block), объединяет результаты.
    Возвращает слова с координатами в ИСХОДНЫХ пикселях кропа.
    """
    if _HAS_OCR_CACHE and assets_dir:
        cached = _load_ocr_cache(crop_img, 'v1', assets_dir)
        if cached is not None:
            return cached
    variants = _make_gray_variants(crop_img)

    # ── Phase 56: Confidence-Weighted Multi-Variant Voting ────────────────────
    # Collect per-variant word lists, then vote instead of plain merge.
    variant_word_lists: list[list[OcrWord]] = []

    # PSM 6 (uniform block) — приоритетно для плотных текстовых сеток (таблицы, штамп-блоки)
    if _has_dense_text_grid(crop_img):
        dense_words = _run_tesseract(variants[0][0], psm=6, up_scale=variants[0][1])
        if dense_words:
            variant_word_lists.append(dense_words)

    # PSM 11 (sparse text) — для разбросанных аннотаций на чертежах
    # PSM 6 (uniform block) добавляем только для первого варианта (grayscale) —
    # подхватывает сплошные блоки текста в таблицах
    # Вариант 3 (equalize) — только PSM 11, без PSM 6 чтобы не удвоить шум
    for vi, (img_variant, up_scale) in enumerate(variants):
        batch11 = _run_tesseract(img_variant, psm=11, up_scale=up_scale)
        if batch11:
            variant_word_lists.append(batch11)
        if vi == 0:
            # PSM 6 на grayscale варианте — дополнительно для табличного текста
            batch6 = _run_tesseract(img_variant, psm=6, up_scale=up_scale)
            if batch6:
                variant_word_lists.append(batch6)

    # PSM 7 (single text line) — для изолированных числовых значений на чистом фоне
    # (слайды 8, 18). Принимаем только слова-измерения/_TECH_MARKER чтобы не добавлять шум.
    batch7 = _run_tesseract(variants[0][0], psm=7, up_scale=variants[0][1])
    batch7_filtered = [
        w for w in batch7
        if _MEASUREMENT.match(w.text.strip()) or _TECH_MARKER.match(w.text.strip())
    ]
    if batch7_filtered:
        variant_word_lists.append(batch7_filtered)

    # Vote across all preprocessing variants
    all_words: list[OcrWord] = _vote_ocr_results(variant_word_lists)

    up_scale_main = variants[0][1] if variants else 2.0
    rotated_words = _try_rotated_ocr(crop_img, up_scale_main)
    all_words = _merge_word_lists(all_words, rotated_words, tol=20)

    # 4x upscale re-OCR для маленьких низкоуверенных регионов (char height < 20px, conf < 0.60)
    small_regions = _find_small_regions(all_words)
    iw, ih = crop_img.size
    for rx0, ry0, rx1, ry1 in small_regions:
        rx0 = max(0, rx0); ry0 = max(0, ry0); rx1 = min(iw, rx1); ry1 = min(ih, ry1)
        if rx1 <= rx0 or ry1 <= ry0:
            continue
        sub = crop_img.crop((rx0, ry0, rx1, ry1))
        from PIL import ImageEnhance as _IE
        sub4x = sub.resize((sub.width * 4, sub.height * 4), Image.LANCZOS).convert("RGB")
        sub4x_enhanced = _IE.Contrast(sub4x.convert("L")).enhance(2.5).convert("RGB")
        sub_words_raw = _run_tesseract(sub4x_enhanced, psm=11, up_scale=4.0)
        # Конвертируем координаты из sub4x-пространства обратно в пространство кропа
        sub_words = [
            OcrWord(w.text, rx0 + w.px0, ry0 + w.py0, rx0 + w.px1, ry0 + w.py1, w.conf)
            for w in sub_words_raw
        ]
        all_words = _merge_word_lists(all_words, sub_words, tol=6)

    # PSM 8 (single-word) recovery pass for small callout text in grid cells with no OCR hits
    import numpy as _np
    _GRID_ROWS, _GRID_COLS = 12, 4
    cell_w = iw / _GRID_COLS
    cell_h = ih / _GRID_ROWS
    # Build set of cells that already contain at least one word center
    occupied: set[tuple[int, int]] = set()
    for w in all_words:
        cx = (w.px0 + w.px1) / 2
        cy = (w.py0 + w.py1) / 2
        col = int(cx / cell_w)
        row = int(cy / cell_h)
        occupied.add((row, col))
    # Find empty-but-dark cells
    gray_arr = _np.array(crop_img.convert("L"))
    psm8_candidates: list[tuple[int, int, int, int]] = []
    for row in range(_GRID_ROWS):
        for col in range(_GRID_COLS):
            if (row, col) in occupied:
                continue
            cx0 = int(col * cell_w)
            cy0 = int(row * cell_h)
            cx1 = min(iw, int((col + 1) * cell_w))
            cy1 = min(ih, int((row + 1) * cell_h))
            if cx1 <= cx0 or cy1 <= cy0:
                continue
            cell_arr = gray_arr[cy0:cy1, cx0:cx1]
            dark_frac = float((cell_arr < 100).sum()) / cell_arr.size
            if dark_frac > 0.05:
                psm8_candidates.append((cx0, cy0, cx1, cy1))
    if 0 < len(psm8_candidates) <= 12:
        for cx0, cy0, cx1, cy1 in psm8_candidates:
            cell_crop = crop_img.crop((cx0, cy0, cx1, cy1))
            cell3x = cell_crop.resize((cell_crop.width * 3, cell_crop.height * 3), Image.LANCZOS).convert("RGB")
            psm8_words_raw = _run_tesseract(cell3x, psm=8, up_scale=3.0)
            psm8_words = [
                OcrWord(w.text, cx0 + w.px0, cy0 + w.py0, cx0 + w.px1, cy0 + w.py1, w.conf)
                for w in psm8_words_raw
                if w.conf > 0.55 and _is_valid_ocr(w.text) and len(w.text) > 3
            ]
            if psm8_words:
                all_words = _merge_word_lists(all_words, psm8_words, tol=10)

    all_words = _reread_low_conf_measurements(all_words, crop_img)
    all_words = _reread_uncertain_floats(all_words, crop_img)

    # Phase 55: Table Structure Recovery — run before returning to merge
    # per-cell OCR results for images that contain grid tables.
    table_cells = _detect_table_cells(crop_img)
    if table_cells:
        print(f"    [Phase 55] Table detected: {len(table_cells)} cells — running per-cell OCR")
        cell_words = _ocr_table_cells(crop_img, table_cells)
        if cell_words:
            all_words = _merge_word_lists(all_words, cell_words, tol=8)
            print(f"    [Phase 55] Cell OCR added {len(cell_words)} words")

    all_words = _dedup_words(all_words)

    if _HAS_OCR_CACHE and assets_dir:
        _save_ocr_cache(crop_img, 'v1', all_words, assets_dir)
    return all_words


# ── Шрифт ────────────────────────────────────────────────────────────────────

def _word_font(name: str | None) -> str:
    if not name:
        return "Times New Roman"
    clean = name.split("+")[-1]
    for suffix in ["-Bold", "-Italic", "-BoldItalic", ",Bold", ",Italic"]:
        clean = clean.replace(suffix, "")
    _MAP = {
        "TimesNewRoman": "Times New Roman",
        "ArialMT":       "Arial",
        "Arial-BoldMT":  "Arial",
    }
    return _MAP.get(clean.strip(), "Times New Roman")


def _rgb_from_fitz(color_int) -> RGBColor:
    if color_int is None:
        return RGBColor(0, 0, 0)
    if isinstance(color_int, float):
        v = int(color_int * 255)
        return RGBColor(v, v, v)
    if isinstance(color_int, (list, tuple)):
        if len(color_int) == 3:
            r, g, b = color_int
            return RGBColor(int(r*255), int(g*255), int(b*255))
        return RGBColor(0, 0, 0)
    r = (color_int >> 16) & 0xFF
    g = (color_int >>  8) & 0xFF
    b =  color_int        & 0xFF
    return RGBColor(r, g, b)


# ── Text box helpers ──────────────────────────────────────────────────────────

def _set_no_line(txBox) -> None:
    from lxml import etree
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    spPr = txBox._element.find(f".//{{{ns}}}spPr")
    if spPr is None:
        return
    ln = spPr.find(f"{{{ns}}}ln")
    if ln is None:
        ln = etree.SubElement(spPr, f"{{{ns}}}ln")
    if ln.find(f"{{{ns}}}noFill") is None:
        etree.SubElement(ln, f"{{{ns}}}noFill")


def _set_no_fill(txBox) -> None:
    """Убираем заливку фона text box (прозрачный)."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    spPr = txBox._element.find(f".//{{{ns_a}}}spPr")
    if spPr is None:
        return
    # Ищем или создаём noFill в spPr
    if spPr.find(f"{{{ns_a}}}noFill") is None:
        # Убираем solidFill если есть
        sf = spPr.find(f"{{{ns_a}}}solidFill")
        if sf is not None:
            spPr.remove(sf)
        etree.SubElement(spPr, f"{{{ns_a}}}noFill")


def _group_lines_by_visual_row(lines: list[dict]) -> list[list[dict]]:
    if not lines:
        return []
    rows: list[list[dict]] = []
    cur_row = [lines[0]]
    cur_y0 = lines[0]["bbox"][1]
    line_h = lines[0]["bbox"][3] - lines[0]["bbox"][1]
    threshold = max(4.0, line_h * 0.6)
    for line in lines[1:]:
        y0 = line["bbox"][1]
        if abs(y0 - cur_y0) <= threshold:
            cur_row.append(line)
        else:
            rows.append(cur_row)
            cur_row = [line]
            cur_y0 = y0
    rows.append(cur_row)
    return rows


def _add_block_textbox(slide, block: dict) -> None:
    bx0, by0, bx1, by1 = block["bbox"]
    w_pt = max(bx1 - bx0, 4.0)
    h_pt = max(by1 - by0, 4.0)
    txBox = slide.shapes.add_textbox(
        Emu(_emu(bx0)), Emu(_emu(by0)),
        Emu(_emu(w_pt)), Emu(_emu(h_pt * 1.15)),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _set_no_line(txBox)
    visual_rows = _group_lines_by_visual_row(block.get("lines", []))
    for ri, row_lines in enumerate(visual_rows):
        p = tf.paragraphs[0] if ri == 0 else tf.add_paragraph()
        for line in sorted(row_lines, key=lambda l: l["bbox"][0]):
            for span in line.get("spans", []):
                txt = span.get("text", "")
                if not txt:
                    continue
                run = p.add_run()
                run.text = txt
                fnt = run.font
                size_pt = span.get("size", 11.0) or 11.0
                fnt.size   = Pt(max(4, size_pt))
                fnt.name   = _word_font(span.get("font"))
                fnt.bold   = bool(span.get("flags", 0) & (1 << 4))
                fnt.italic = bool(span.get("flags", 0) & (1 << 1))
                fnt.color.rgb = _rgb_from_fitz(span.get("color"))


def _merge_close_lines(lines: list[list[OcrWord]], img_native_w: int) -> list[list[OcrWord]]:
    """Merges vertically close OCR lines into single multi-line groups.
    Улучшает поиск и логическую группировку аннотаций типа «СТАЛЬ / 45 / ГОСТ 1050».
    """
    from statistics import mean as _mean

    def _avg_y(ln: list[OcrWord]) -> float:
        return sum((w.py0 + w.py1) / 2 for w in ln) / len(ln)

    def _avg_h_words(a: list[OcrWord], b: list[OcrWord]) -> float:
        all_h = [w.py1 - w.py0 for w in a + b if w.py1 - w.py0 > 0]
        return _mean(all_h) if all_h else 1.0

    def _x_range(ln: list[OcrWord]) -> tuple[int, int]:
        return (min(w.px0 for w in ln), max(w.px1 for w in ln))

    for _pass in range(3):
        merged = False
        new_lines: list[list[OcrWord]] = []
        i = 0
        while i < len(lines):
            if i + 1 < len(lines):
                a = lines[i]
                b = lines[i + 1]
                # Only merge when b is strictly below a
                if _avg_y(b) > _avg_y(a):
                    avg_h = _avg_h_words(a, b)
                    gap = min(w.py0 for w in b) - max(w.py1 for w in a)
                    ax0, ax1 = _x_range(a)
                    bx0, bx1 = _x_range(b)
                    overlap = max(0, min(ax1, bx1) - max(ax0, bx0))
                    min_span = max(1, min(ax1 - ax0, bx1 - bx0))
                    overlap_frac = overlap / min_span
                    combined_w = max(ax1, bx1) - min(ax0, bx0)
                    width_ok = (img_native_w <= 0) or (combined_w < img_native_w * 0.75)
                    if gap < avg_h * 1.8 and overlap_frac > 0.25 and width_ok:
                        combined = sorted(a + b, key=lambda w: (w.py0, w.px0))
                        new_lines.append(combined)
                        i += 2
                        merged = True
                        continue
            new_lines.append(lines[i])
            i += 1
        lines = new_lines
        if not merged:
            break
    return lines


def _group_ocr_into_lines(words: list[OcrWord],
                          img_native_w: int = 0) -> list[list[OcrWord]]:
    """Группирует OCR-слова с близкими y-координатами в визуальные строки.
    Phase 65: Adaptive Line Gap — порог группировки вычисляется как median_h * 0.6,
    где median_h — медианная высота слова. Это адаптивнее фиксированного avg_h * 0.5
    и устойчивее к выбросам (крупные заголовки не растягивают порог).
    Удаляет строки шириной > 80% изображения.
    """
    if not words:
        return []

    # Phase 65: вычисляем медианную высоту слова один раз перед группировкой
    heights = [w.py1 - w.py0 for w in words if w.py1 - w.py0 > 0]
    if not heights:
        median_h = 12
    else:
        heights.sort()
        median_h = heights[len(heights) // 2]
    median_h = max(median_h, 4)  # минимальная защита от вырождения

    # Порог Y-расстояния для слов на одной строке
    line_gap_threshold = median_h * 0.6
    # Порог Y-расстояния для разделения отдельных текстовых блоков
    block_split_threshold = median_h * 1.5

    # Сортируем по центру y, а не по py0 — стабильнее при разных baseline
    sorted_w = sorted(words, key=lambda w: (w.py0 + w.py1) / 2)
    lines: list[list[OcrWord]] = []
    cur_line = [sorted_w[0]]
    cur_cy = (sorted_w[0].py0 + sorted_w[0].py1) / 2
    for w in sorted_w[1:]:
        cy = (w.py0 + w.py1) / 2
        if abs(cy - cur_cy) <= line_gap_threshold:
            cur_line.append(w)
            # Обновляем центр строки как среднее всех слов
            cur_cy = sum((ww.py0 + ww.py1) / 2 for ww in cur_line) / len(cur_line)
        else:
            lines.append(sorted(cur_line, key=lambda ww: ww.px0))
            cur_line = [w]
            cur_cy = cy
    lines.append(sorted(cur_line, key=lambda ww: ww.px0))

    # Фильтруем строки шире 80% изображения (без PSM 6 таких быть не должно)
    if img_native_w > 0:
        max_line_w = img_native_w * 0.80
        lines = [
            ln for ln in lines
            if (max(w.px1 for w in ln) - min(w.px0 for w in ln)) <= max_line_w
        ]

    # Фильтруем "штриховочные" строки: >70% слов ≤3 символа И средний conf < 0.55
    _CAPS_NOISE_PAT = _re.compile(r'^[A-Z]{2,4}[!;.,]?$')

    def _is_hatch_line(ln: list[OcrWord]) -> bool:
        short = sum(1 for w in ln if len(w.text.strip()) <= 3)
        avg_conf = sum(w.conf for w in ln) / len(ln)
        return (short / len(ln)) > 0.70 and avg_conf < 0.55

    def _is_caps_noise_line(ln: list[OcrWord]) -> bool:
        """Строки где >55% слов — 2-4 char ALL CAPS Latin без цифр и кириллицы."""
        caps_noise = sum(
            1 for w in ln
            if _CAPS_NOISE_PAT.match(w.text.strip())
            and not _CYRILLIC.search(w.text)
            and not _re.search(r"\d", w.text)
        )
        avg_conf = sum(w.conf for w in ln) / len(ln)
        return (caps_noise / len(ln)) > 0.55 and avg_conf < 0.65

    def _is_dense_noise_line(ln: list[OcrWord]) -> bool:
        """Phase 62: строки с высоким повторением символов — штриховочный мусор."""
        joined = "".join(w.text for w in ln)
        if len(joined) < 4:
            return False
        char_counts: dict[str, int] = {}
        for c in joined:
            char_counts[c] = char_counts.get(c, 0) + 1
        repeated = sum(v for v in char_counts.values() if v > 1)
        rep_ratio = repeated / len(joined)
        has_digit = bool(_re.search(r"\d", joined))
        non_alnum = sum(1 for c in joined if not _ALPHANUM.match(c))
        non_alnum_ratio = non_alnum / len(joined)
        if rep_ratio > 0.65 and not has_digit:
            return True
        if non_alnum_ratio > 0.5:
            return True
        return False

    lines = [ln for ln in lines if not _is_hatch_line(ln) and not _is_caps_noise_line(ln) and not _is_dense_noise_line(ln)]

    # Разбиваем "широкие" строки по горизонтальным разрывам.
    # Если расстояние между соседними словами > 4× средней ширины слова,
    # это скорее всего две разные аннотации на одной y-высоте.
    split_lines: list[list[OcrWord]] = []
    for ln in lines:
        if len(ln) <= 1:
            split_lines.append(ln)
            continue
        avg_w = sum(w.px1 - w.px0 for w in ln) / len(ln)
        gap_thresh = max(avg_w * 4, 60)
        cur = [ln[0]]
        for w in ln[1:]:
            gap = w.px0 - cur[-1].px1
            if gap > gap_thresh:
                split_lines.append(cur)
                cur = [w]
            else:
                cur.append(w)
        split_lines.append(cur)
    lines = split_lines

    # Merge-back step: if two consecutive gap-split lines have the same
    # approximate y-center (within 1 line-height) and their combined width
    # still fits within 80% of image width, they are likely one annotation
    # that was falsely split.
    merged: list[list[OcrWord]] = []
    i = 0
    while i < len(lines):
        ln = lines[i]
        if i + 1 < len(lines):
            nxt = lines[i + 1]
            # Average y-centers
            cy1 = sum((w.py0 + w.py1) / 2 for w in ln) / len(ln)
            cy2 = sum((w.py0 + w.py1) / 2 for w in nxt) / len(nxt)
            avg_h_ln = sum(w.py1 - w.py0 for w in ln) / len(ln)
            avg_h_nxt = sum(w.py1 - w.py0 for w in nxt) / len(nxt)
            line_h = (avg_h_ln + avg_h_nxt) / 2
            # Combined bounding box width
            all_words = ln + nxt
            combined_w = max(w.px1 for w in all_words) - min(w.px0 for w in all_words)
            width_ok = (img_native_w <= 0) or (combined_w <= img_native_w * 0.80)
            if abs(cy1 - cy2) <= line_h and width_ok:
                merged.append(sorted(all_words, key=lambda ww: ww.px0))
                i += 2
                continue
        merged.append(ln)
        i += 1
    lines = merged

    return lines


_MEAS_TRAIL   = _re.compile(r'^([—\-±=]?\d+[.,]\d+)[a-zA-Z!°]+$')
_LEAD_NOISE   = _re.compile(r'^[\\~#*°]([=\-+]?\d)')
_BRACKET_DIG  = _re.compile(r'^\[(\d)')
_EQ_LETTER    = _re.compile(r'^=[a-zA-Z]$')
# Trailing comma/period from measurement values: -0.300, → -0.300, -0,88. → -0,88
_MEAS_TRAIL_PUNCT = _re.compile(r'^([—\-±=]?\d+[.,]\d+)[.,]+$')
# Trailing punct from digit-only tokens: 10300! → 10300
_DIGIT_TRAIL  = _re.compile(r'^(\d[\d.,]*)[\s!;.°]+$')
# 2-3 символьные ALL CAPS Latin без цифр/кириллицы — штриховочный шум
_CAPS2        = _re.compile(r'^[A-Z]{2,3}[!;.,]?$')
_INIT_CAP_NOISE = _re.compile(r'^[A-Z][a-z]{1,3}$')


def _normalize_measurement(text: str) -> str:
    """Phase 69: нормализует форматы измерений — запятая→точка, DN с пробелом, trailing zeros."""
    t = text.strip()
    t = _re.sub(r'(\d),(\d)', r'\1.\2', t)
    t = _re.sub(r'(\d\.\d*[1-9])0+$', r'\1', t)
    t = _re.sub(r'\b(DN|Dn|dn)\s+(\d+)', r'DN\2', t)
    t = _re.sub(r'^(-?\d+\.\d+)[a-eg-z]$', r'\1', t)
    return t


def _clean_token(t: str) -> str:
    """Очищает OCR-токен от мусорных prefix/suffix: trailing буква, leading \\~*°[ перед числом."""
    # trailing single letter/symbol after decimal: -0.07g → -0.07
    m = _MEAS_TRAIL.match(t)
    if m:
        return m.group(1)
    # leading noise char before digit: \=0.36 → =0.36, ~0.30 → 0.30, *0.20 → 0.20, °0.07q → 0.07q
    if _LEAD_NOISE.match(t):
        t2 = t[1:]
        # second pass: handle compound cases (°0.07q → 0.07q → 0.07)
        m2 = _MEAS_TRAIL.match(t2)
        if m2:
            return m2.group(1)
        return t2
    if _BRACKET_DIG.match(t):
        return t[1:]
    # trailing comma/period from measurement values: -0.300, → -0.300
    m = _MEAS_TRAIL_PUNCT.match(t)
    if m:
        return m.group(1)
    # trailing punct from digit tokens: 10300! → 10300
    m = _DIGIT_TRAIL.match(t)
    if m:
        return m.group(1)
    if _MEASUREMENT.match(t) or _FLOAT_2DEC.match(t):
        t = _normalize_measurement(t)
    return t


def _min_conf_for_line(words: list) -> float:
    """Returns the minimum confidence threshold (0-1 float) for a line of OcrWords.

    Rules (in priority order):
    - Any word matches _MEASUREMENT -> threshold 0.20  (measurement values on hatched backgrounds)
    - Any word contains Cyrillic    -> threshold 0.22  (Cyrillic OCR is more reliable)
    - Default                       -> threshold 0.28
    """
    has_measurement = any(_MEASUREMENT.match(w.text.strip()) for w in words)
    if has_measurement:
        return 0.20
    has_cyrillic = any(_CYRILLIC.search(w.text) for w in words)
    if has_cyrillic:
        return 0.22
    return 0.28


def _iou(a: tuple, b: tuple) -> float:
    """Compute Intersection over Union for two (px0,py0,px1,py1) tuples."""
    ix0 = max(a[0], b[0]); iy0 = max(a[1], b[1])
    ix1 = min(a[2], b[2]); iy1 = min(a[3], b[3])
    iw = max(0, ix1 - ix0); ih = max(0, iy1 - iy0)
    inter = iw * ih
    if inter == 0:
        return 0.0
    area_a = (a[2] - a[0]) * (a[3] - a[1])
    area_b = (b[2] - b[0]) * (b[3] - b[1])
    union = area_a + area_b - inter
    return inter / union if union > 0 else 0.0


def _dedup_words(words: list) -> list:
    """Remove duplicate OcrWords with IoU > 0.5, keep higher-conf one."""
    kept: list = []
    for w in words:
        bbox_w = (w.px0, w.py0, w.px1, w.py1) if hasattr(w, 'px0') else (w[1], w[2], w[3], w[4])
        conf_w = w.conf if hasattr(w, 'conf') else 50
        duplicate = False
        for i, k in enumerate(kept):
            bbox_k = (k.px0, k.py0, k.px1, k.py1) if hasattr(k, 'px0') else (k[1], k[2], k[3], k[4])
            conf_k = k.conf if hasattr(k, 'conf') else 50
            if _iou(bbox_w, bbox_k) > 0.5:
                if conf_w > conf_k:
                    kept[i] = w
                duplicate = True
                break
        if not duplicate:
            kept.append(w)
    return kept


def _overlaps_native(
    ocr_bbox_pt: tuple[float, float, float, float],
    native_bboxes_pt: list[tuple[float, float, float, float]],
    tol: float = 5.0,
) -> bool:
    """Return True if the OCR line center falls inside any native text bbox (expanded by tol pt)."""
    cx = (ocr_bbox_pt[0] + ocr_bbox_pt[2]) / 2
    cy = (ocr_bbox_pt[1] + ocr_bbox_pt[3]) / 2
    for x0, y0, x1, y1 in native_bboxes_pt:
        if (x0 - tol) <= cx <= (x1 + tol) and (y0 - tol) <= cy <= (y1 + tol):
            return True
    return False


def _set_textbox_font(tf, font_size_pt: float) -> None:
    """OCR text boxes: тёмный читаемый текст без проверки орфографии."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(max(6, min(font_size_pt, 72)))
            run.font.color.rgb = RGBColor(0x15, 0x15, 0x15)  # почти чёрный
            # Отключаем проверку орфографии для OCR-текста
            rPr = run._r.get_or_add_rPr()
            rPr.set('noProof', '1')
            rPr.set('lang', 'ru-RU')


def _add_ocr_line_textbox(
    slide,
    line_words: list[OcrWord],
    img_pt_x: float,
    img_pt_y: float,
    img_pt_w: float,
    img_pt_h: float,
    crop_px_w: int,
    crop_px_h: int,
) -> None:
    """Добавляет один textbox для строки OCR-слов (объединённых по y-координате)."""
    if crop_px_w <= 0 or crop_px_h <= 0 or not line_words:
        return
    sx = img_pt_w / crop_px_w
    sy = img_pt_h / crop_px_h

    px0 = min(w.px0 for w in line_words)
    py0 = min(w.py0 for w in line_words)
    px1 = max(w.px1 for w in line_words)
    py1 = max(w.py1 for w in line_words)

    x_pt = img_pt_x + px0 * sx
    y_pt = img_pt_y + py0 * sy
    w_pt = max((px1 - px0) * sx * 1.10, 4.0)
    h_pt = max((py1 - py0) * sy * 1.4, 4.0)
    # Шрифт: классификация по средней высоте слов в пунктах (Phase 45)
    word_heights = sorted((w.py1 - w.py0) for w in line_words)
    p25_h = word_heights[len(word_heights) // 4]
    import statistics as _statistics
    mean_h_px = _statistics.mean((w.py1 - w.py0) for w in line_words)
    mean_h_pt = mean_h_px * sy
    # Высококонфидентная строка (conf > 0.85 у всех слов) — cap +2pt
    all_high_conf = all(w.conf > 0.85 for w in line_words)
    conf_bonus = 2.0 if all_high_conf else 0.0
    if mean_h_pt > 18:
        # Крупный текст — заголовки/шапки
        font_size = min(28.0 + conf_bonus, mean_h_pt * 0.75)
    elif mean_h_pt > 10:
        # Нормальный текст
        font_size = min(14.0 + conf_bonus, max(5.0, p25_h * sy * 0.72))
    else:
        # Мелкий текст — выноски, сноски
        font_size = min(10.0 + conf_bonus, max(5.0, mean_h_pt * 0.80))

    txBox = slide.shapes.add_textbox(
        Emu(_emu(x_pt)), Emu(_emu(y_pt)),
        Emu(_emu(w_pt)), Emu(_emu(h_pt)),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _set_no_line(txBox)
    _set_no_fill(txBox)

    p = tf.paragraphs[0]
    run = p.add_run()
    # Inline-фильтры шумовых токенов
    _inline_noise   = _re.compile(r'^["\'\`\*\\\|\^~<>{}\[\]@#!]+$|^["\'\`\*][a-z]{1,3}$')
    _inline_low_lat = _re.compile(r'^[a-z]{1,5}[).,!|]*$')

    def _keep_word(w: OcrWord) -> tuple[bool, str]:
        t = w.text.strip()
        # Пунктуационный или цитатный мусор
        if _inline_noise.match(t):
            return False, t
        # Строчный латинский без цифр/кириллицы
        if _inline_low_lat.match(t) and not _CYRILLIC.search(t) and not _re.search(r"\d", t):
            return False, t
        # ≥3-символьные ALL CAPS Latin без кириллицы (уже не проходят _is_valid_ocr)
        if _ALL_CAPS_LATIN.match(t) and not _CYRILLIC.search(t):
            return False, t
        # 2-3-символьные ALL CAPS Latin — штриховочный шум (SS, CA, LN, FS, Ne→No)
        if _CAPS2.match(t) and not _CYRILLIC.search(t) and not _re.search(r"\d", t):
            return False, t
        # Начальная заглавная + 1-3 строчные Latin без цифр/кириллицы — штриховочный шум (Ne, Kr, Go, Gy, Ren)
        if _INIT_CAP_NOISE.match(t) and not _CYRILLIC.search(t) and not _re.search(r"\d", t) and w.conf < 0.75:
            return False, t
        # = + одна буква: =f, =g — шум
        if _EQ_LETTER.match(t):
            return False, t
        # Очищаем токен: trailing буква после числа, leading \~[ перед числом
        cleaned = _clean_token(t)
        if _DASH_SEQ.match(cleaned) and not _CYRILLIC.search(cleaned):
            return False, cleaned
        return True, cleaned

    cleaned_parts: list[str] = []
    for w in line_words:
        keep, cleaned = _keep_word(w)
        if keep and cleaned:
            cleaned_parts.append(cleaned)

    if not cleaned_parts:
        cleaned_parts = [w.text for w in line_words]
    run.text = " ".join(cleaned_parts)
    run.font.size  = Pt(font_size)
    run.font.name  = "Arial"
    run.font.color.rgb = RGBColor(0x15, 0x15, 0x15)
    # Отключаем орфографию на главном run
    from pptx.oxml.ns import qn as _qn
    _rPr = run._r.get_or_add_rPr()
    _rPr.set('noProof', '1')
    _rPr.set('lang', 'ru-RU')
    _set_textbox_font(tf, font_size)


# ── Заголовок слайда (outline panel) ─────────────────────────────────────────

def _infer_slide_title(words: list[str], slide_idx: int) -> str:
    """Infers a slide title for PowerPoint outline panel.

    Strategy:
    1. Longest ALL_CAPS run of >= 2 consecutive words.
    2. First 6 words when total words >= 3.
    3. Fallback: "Слайд {slide_idx + 1}".
    Result truncated to 60 chars.
    """
    _MAX = 60
    _ALLCAPS = _re.compile(r'^[А-ЯЁA-Z0-9\-/:.,%°]+$')

    def _truncate(s: str) -> str:
        return s[:_MAX].rstrip() if len(s) > _MAX else s

    if not words:
        return f"Слайд {slide_idx + 1}"

    # Strategy 1: longest run of ALL_CAPS words
    best_caps: list[str] = []
    cur_caps: list[str] = []
    for w in words:
        if _ALLCAPS.match(w.strip()):
            cur_caps.append(w.strip())
        else:
            if len(cur_caps) > len(best_caps):
                best_caps = cur_caps
            cur_caps = []
    if len(cur_caps) > len(best_caps):
        best_caps = cur_caps
    if len(best_caps) >= 2:
        return _truncate(" ".join(best_caps))

    # Strategy 2: first run of >= 3 words (any case)
    if len(words) >= 3:
        return _truncate(" ".join(words[:6]))

    return f"Слайд {slide_idx + 1}"


def _set_slide_title(slide, title: str) -> None:
    """Sets slide title text for PowerPoint outline panel.

    Tries the title placeholder (idx=0) first. Falls back to a tiny 1pt
    white invisible textbox tagged as a title placeholder via XML.
    """
    # Try existing title placeholder
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            try:
                ph.text = title
                return
            except Exception:
                break

    # Fallback: tiny invisible textbox at origin
    txBox = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1), Emu(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(1)
    run.font.color.rgb = RGBColor(255, 255, 255)
    _set_no_line(txBox)

    # Tag as title placeholder so PowerPoint outline panel picks it up
    from lxml import etree as _etree
    _PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    try:
        _nvSpPr = txBox._element[0]  # nvSpPr
        _nvPr = _nvSpPr.find(f"{{{_PML_NS}}}nvPr")
        if _nvPr is None:
            _nvPr = _etree.SubElement(_nvSpPr, f"{{{_PML_NS}}}nvPr")
        if _nvPr.find(f"{{{_PML_NS}}}ph") is None:
            ph_el = _etree.SubElement(_nvPr, f"{{{_PML_NS}}}ph")
            ph_el.set("type", "title")
    except Exception:
        pass  # Best-effort — outline panel degrades gracefully


# ── Фон слайда ───────────────────────────────────────────────────────────────

def _set_slide_bg(slide, rgb: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _has_vector_content(
    span_rects: list[tuple], img_rects: list[tuple], pw: float, ph: float
) -> bool:
    text_area = sum((x1-x0)*(y1-y0) for x0,y0,x1,y1 in span_rects)
    img_area  = sum((x1-x0)*(y1-y0) for x0,y0,x1,y1 in img_rects)
    covered   = (text_area + img_area) / (pw * ph) if pw * ph > 0 else 0
    return covered < 0.95


# ── Главная функция ───────────────────────────────────────────────────────────

def emit_pptx_slides(
    ir,
    output_path: Path,
    pdf_path: Path,
    *,
    dpi: int = 200,
    assets_dir: Path | None = None,
    ocr_images: bool = True,
) -> Path:
    """
    Конвертирует slide-PDF в PPTX.
    Каждый слайд: белый фон + вектор-слой + отдельные изображения (с OCR) + text boxes.
    """
    _search_index: list[list[str]] = []

    work_dir = assets_dir or output_path.parent / "pptx_assets"
    raw_dir  = work_dir / "raw"
    vec_dir  = work_dir / "vector"
    img_dir  = work_dir / "images"
    msk_dir  = work_dir / "img_masked"
    for d in (raw_dir, vec_dir, img_dir, msk_dir):
        d.mkdir(parents=True, exist_ok=True)

    scale = dpi / 72.0
    mat   = fitz.Matrix(scale, scale)

    pdf_doc   = fitz.open(pdf_path)
    page_data = []

    for i in range(pdf_doc.page_count):
        fz_page = pdf_doc.load_page(i)
        rect    = fz_page.rect
        pw, ph  = rect.width, rect.height

        # ── 1. Рендер страницы без текстового слоя ───────────────────────
        # Используем временную копию страницы с redact-аннотациями для текста,
        # чтобы фоновое изображение не содержало нативный текст (он будет
        # вставлен как редактируемые text boxes отдельно).
        try:
            _tmp_doc = fitz.open()
            _tmp_doc.insert_pdf(pdf_doc, from_page=i, to_page=i)
            _tmp_page = _tmp_doc[0]
            # Redact all text spans from temporary copy
            _txt_dict_tmp = _tmp_page.get_text("dict", flags=0)
            for _blk in _txt_dict_tmp.get("blocks", []):
                if _blk.get("type") == 0:
                    for _ln in _blk.get("lines", []):
                        for _sp in _ln.get("spans", []):
                            if _sp.get("text", "").strip():
                                _tmp_page.add_redact_annot(fitz.Rect(_sp["bbox"]))
            _tmp_page.apply_redactions(
                images=fitz.PDF_REDACT_IMAGE_NONE,
                graphics=fitz.PDF_REDACT_LINE_ART_NONE,
            )
            pix = _tmp_page.get_pixmap(matrix=mat, alpha=False)
            _tmp_doc.close()
        except Exception:
            # Fallback: render with text (редактируемость теряется, но файл создаётся)
            pix = fz_page.get_pixmap(matrix=mat, alpha=False)
        raw_png = raw_dir / f"slide_{i:03d}.png"
        pix.save(str(raw_png))
        img_full = Image.open(raw_png).convert("RGB")
        iw, ih   = img_full.size

        # ── 2. Текстовые блоки (native PDF text) ──────────────────────────
        text_dict   = fz_page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        text_blocks = [b for b in text_dict.get("blocks", []) if b.get("type") == 0]

        span_rects: list[tuple] = []
        for block in text_blocks:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    if span.get("text", "").strip():
                        span_rects.append(tuple(span["bbox"]))

        # ── 3. Изображения (get_image_info с xrefs) ───────────────────────
        img_infos = fz_page.get_image_info(xrefs=True)
        img_rects: list[tuple] = [
            tuple(info["bbox"]) for info in img_infos
            if info.get("bbox") and (info["bbox"][2]-info["bbox"][0]) > 4
        ]

        # ── 4. Кропы + OCR для каждого изображения ────────────────────────
        cropped_data: list[dict] = []
        for bi, info in enumerate(img_infos):
            bx0, by0, bx1, by1 = info["bbox"]
            pw_img = bx1 - bx0
            ph_img = by1 - by0
            if pw_img < 4 or ph_img < 4:
                continue

            # Пиксельные координаты кропа в рендере (для отображения/маскировки)
            cpx0 = max(0, int(bx0 * scale))
            cpy0 = max(0, int(by0 * scale))
            cpx1 = min(iw, int(bx1 * scale))
            cpy1 = min(ih, int(by1 * scale))
            if cpx1 - cpx0 < 4 or cpy1 - cpy0 < 4:
                continue

            raw_crop = img_full.crop((cpx0, cpy0, cpx1, cpy1))
            crop_w, crop_h = raw_crop.size

            # Пробуем нативное извлечение изображения (выше разрешение → лучше OCR)
            xref = info.get("xref", 0)
            native_img: Image.Image | None = None
            if xref:
                try:
                    import io as _io
                    img_data = pdf_doc.extract_image(xref)
                    raw_bytes = img_data.get("image")
                    if raw_bytes:
                        native_img = Image.open(_io.BytesIO(raw_bytes)).convert("RGB")
                except Exception:
                    native_img = None

            ocr_source = native_img if native_img else raw_crop
            native_w, native_h = ocr_source.size

            # OCR
            ocr_words: list[OcrWord] = []
            # Пропускаем маленькие декоративные изображения (логотипы, иконки)
            # Реальные чертежи всегда >= 900px по обоим измерениям
            _is_tiny = min(native_w, native_h) < 600
            if ocr_images and native_w > 30 and native_h > 20 and not _is_tiny:
                print(f"  Слайд {i}, image {bi}: OCR ({native_w}x{native_h}px, "
                      f"{'native' if native_img else 'rendered'})...")
                ocr_words = _ocr_crop(ocr_source)
                print(f"    → {len(ocr_words)} слов")

            # Если OCR делался на нативном изображении, масштабируем координаты
            # в пиксели рендер-кропа (для маскировки)
            if native_img and ocr_words:
                sx = crop_w / native_w
                sy = crop_h / native_h
                render_words = [
                    OcrWord(w.text,
                            int(w.px0 * sx), int(w.py0 * sy),
                            int(w.px1 * sx), int(w.py1 * sy),
                            w.conf)
                    for w in ocr_words
                ]
            else:
                render_words = ocr_words

            # Текст в кропе НЕ маскируем: OCR textbox накладывается поверх изображения
            # с точным позиционированием и перекрывает растровый текст.
            # Маскировка создавала "дыры" в чертеже — убрана.
            masked_crop = raw_crop
            masked_path = msk_dir / f"slide_{i:03d}_img_{bi}.png"
            masked_crop.save(str(masked_path))

            # OCR-слова хранятся в нативных координатах (для точного позиционирования)
            cropped_data.append({
                "path":      masked_path,
                "bbox_pt":   (bx0, by0, bx1, by1),
                "crop_px":   (native_w, native_h),  # нативный размер для масштабирования
                "words":   ocr_words,
            })

        # ── 5. Вектор-слой (без текста и без image-blocks) ────────────────
        bg_color = _page_bg_color(img_full)
        vec_img  = _mask_rects(
            img_full,
            _pt_rects_to_px(span_rects, scale),
        )
        vec_img  = _mask_rects(
            vec_img,
            _pt_rects_to_px(img_rects, scale),
            fill=bg_color,
        )

        has_vec = _has_vector_content(span_rects, img_rects, pw, ph)
        vec_png: Path | None = None
        if has_vec:
            vec_png = vec_dir / f"slide_{i:03d}_vec.png"
            vec_img.save(str(vec_png))

        page_data.append((bg_color, vec_png, cropped_data, text_blocks, pw, ph))
        print(f"  Слайд {i}: {len(img_infos)} imgs ({sum(len(c['words']) for c in cropped_data)} OCR-слов), {len(text_blocks)} text blocks")

    pdf_doc.close()

    # ── 6. Сборка PPTX ───────────────────────────────────────────────────────
    prs = Presentation()

    for i, (bg_color, vec_png, cropped_data, text_blocks, pw, ph) in enumerate(page_data):
        prs.slide_width  = Emu(_emu(pw))
        prs.slide_height = Emu(_emu(ph))
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        _set_slide_bg(slide, bg_color)

        # Slide title for PowerPoint outline panel
        _title_words: list[str] = []
        for _blk in text_blocks:
            for _ln in _blk.get("lines", []):
                for _sp in _ln.get("spans", []):
                    _title_words.extend(_sp.get("text", "").split())
                    if len(_title_words) >= 20:
                        break
                if len(_title_words) >= 20:
                    break
            if len(_title_words) >= 20:
                break
        _set_slide_title(slide, _infer_slide_title(_title_words, i))

        # Вектор-слой (декор, рамки, линии)
        if vec_png:
            slide.shapes.add_picture(
                str(vec_png), Emu(0), Emu(0), Emu(_emu(pw)), Emu(_emu(ph)),
            )

        # Изображения (замазанные кропы) на своих позициях
        for cd in cropped_data:
            bx0, by0, bx1, by1 = cd["bbox_pt"]
            slide.shapes.add_picture(
                str(cd["path"]),
                Emu(_emu(bx0)), Emu(_emu(by0)),
                Emu(_emu(bx1 - bx0)), Emu(_emu(by1 - by0)),
            )

        # Collect native text block bboxes for dedup check
        native_bboxes_pt = [tuple(block["bbox"]) for block in text_blocks]

        # OCR text boxes поверх каждого изображения (группировка по строкам)
        _slide_words: list[dict] = []
        for cd in cropped_data:
            bx0, by0, bx1, by1 = cd["bbox_pt"]
            cw, ch = cd["crop_px"]
            sx = (bx1 - bx0) / cw if cw > 0 else 1.0
            sy = (by1 - by0) / ch if ch > 0 else 1.0

            # PHASE 44 VALIDATION: проверяем соответствие аспектных соотношений.
            # PDF bbox и нативный размер изображения должны масштабироваться одинаково.
            # Расхождение > 5% означает, что PDF растянул изображение — используем
            # минимальный масштаб (conservative), чтобы текст не вышел за границу.
            if sx > 0 and sy > 0:
                ratio_diff = abs(sx - sy) / max(sx, sy)
                if ratio_diff > 0.05:
                    print(
                        f"  [WARN] aspect mismatch on slide {i} img bbox "
                        f"({bx0:.1f},{by0:.1f},{bx1:.1f},{by1:.1f}): "
                        f"sx={sx:.4f} sy={sy:.4f} diff={ratio_diff:.1%} — "
                        f"using conservative min scale"
                    )
                    s_cons = min(sx, sy)
                    bx1 = bx0 + s_cons * cw
                    by1 = by0 + s_cons * ch
                    sx = s_cons
                    sy = s_cons

            lines = _group_ocr_into_lines(cd["words"], img_native_w=cw)
            lines = _merge_close_lines(lines, cw)
            for line_words in lines:
                if not line_words:
                    continue
                # Phase 52: per-token confidence filtering
                _line_threshold = _min_conf_for_line(line_words)
                line_words = [
                    w for w in line_words
                    if w.conf >= _line_threshold or w.text.strip().lower() in _TECH_WHITELIST
                ]
                if not line_words:
                    continue
                px0 = min(w.px0 for w in line_words)
                py0 = min(w.py0 for w in line_words)
                px1 = max(w.px1 for w in line_words)
                py1 = max(w.py1 for w in line_words)
                ocr_bbox_pt = (
                    bx0 + px0 * sx,
                    by0 + py0 * sy,
                    bx0 + px1 * sx,
                    by0 + py1 * sy,
                )
                if _overlaps_native(ocr_bbox_pt, native_bboxes_pt):
                    continue
                _add_ocr_line_textbox(
                    slide, line_words,
                    img_pt_x=bx0, img_pt_y=by0,
                    img_pt_w=bx1-bx0, img_pt_h=by1-by0,
                    crop_px_w=cw, crop_px_h=ch,
                )
                # Collect OCR line for search index
                _avg_conf = sum(w.conf for w in line_words) / len(line_words)
                _line_text = " ".join(w.text for w in line_words)
                _slide_words.append({
                    "text": _line_text,
                    "x_pt": round(ocr_bbox_pt[0], 2),
                    "y_pt": round(ocr_bbox_pt[1], 2),
                    "conf": round(_avg_conf, 4),
                    "source": "ocr",
                })

        # Native PDF text boxes
        for block in text_blocks:
            _add_block_textbox(slide, block)
            for _line in block.get("lines", []):
                for _span in _line.get("spans", []):
                    _txt = _span.get("text", "").strip()
                    if not _txt:
                        continue
                    _sx0, _sy0, _sx1, _sy1 = _span["bbox"]
                    _slide_words.append({
                        "text": _txt,
                        "x_pt": round(_sx0, 2),
                        "y_pt": round(_sy0, 2),
                        "conf": 1.0,
                        "source": "native",
                    })

        # Build per-slide token list for v2 index.
        # Lines are stored as sentinel-prefixed entries so _build_search_index
        # can reconstruct both individual words and full line strings.
        _slide_tokens: list[str] = []
        for _entry in _slide_words:
            _line_txt = _entry["text"]
            # Sentinel prefix '\x00' marks a full line entry
            _slide_tokens.append("\x00" + _line_txt)
        _search_index.append(_slide_tokens)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    _index_path = output_path.with_suffix(".search_index.json")
    _index_v2 = _build_search_index(_search_index)
    _index_v2["source_pdf"] = pdf_path.name
    _index_path.write_text(
        json.dumps(_index_v2, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"  Search index v2 → {_index_path}")

    return output_path


def _build_search_index(slide_texts: list[list[str]]) -> dict:
    """Build v2 search index from per-slide word lists.

    Args:
        slide_texts: list of lists; each inner list contains individual token
                     strings for one slide (already split by space from line text).

    Returns:
        v2 index dict with version, slides, and global_index.
    """
    slides_out: list[dict] = []
    global_index: dict[str, list[int]] = {}

    for slide_idx, tokens in enumerate(slide_texts):
        words: list[str] = []
        lines: list[str] = []
        measurements: list[str] = []
        keywords: list[str] = []

        _tok_buf: list[str] = []
        for tok in tokens:
            if tok.startswith("\x00"):
                line_text = tok[1:]
                lines.append(line_text)
                _tok_buf.extend(line_text.split())
            else:
                _tok_buf.append(tok)

        for w in _tok_buf:
            w_strip = w.strip()
            if not w_strip:
                continue
            words.append(w_strip)
            if _MEASUREMENT.match(w_strip) or _TECH_MIXED.match(w_strip):
                if w_strip not in measurements:
                    measurements.append(w_strip)
            if (
                w_strip.lower() in _CYRILLIC_ENGINEERING
                or (len(w_strip) >= 6 and bool(_CYRILLIC.fullmatch(w_strip)))
            ):
                if w_strip not in keywords:
                    keywords.append(w_strip)
            key = w_strip.lower()
            if key not in global_index:
                global_index[key] = []
            if slide_idx not in global_index[key]:
                global_index[key].append(slide_idx)

        slides_out.append({
            "slide_idx": slide_idx,
            "word_count": len(words),
            "words": words,
            "lines": lines,
            "measurements": measurements,
            "keywords": keywords,
        })

    return {
        "version": 2,
        "slides": slides_out,
        "global_index": global_index,
    }
